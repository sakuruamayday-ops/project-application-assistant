from __future__ import annotations

import argparse
import csv
import hashlib
import html
import json
import os
import plistlib
import re
import shutil
import sqlite3
import subprocess
import sys
import tempfile
import unicodedata
import zipfile
from collections import Counter
from datetime import datetime, timezone
from pathlib import Path


MAX_TEXT_CHARS = 200_000
MIN_TEXT_CHARS = 40
DEFAULT_CACHE_PATH = Path(
    os.environ.get(
        "JIAOTANG_EXTRACTION_CACHE",
        Path.home() / ".cache/project-application-assistant/knowledge-extraction-cache.jsonl",
    )
)
LEGACY_CONVERSION_ROOT = Path(
    os.environ.get(
        "JIAOTANG_LEGACY_CONVERSION_ROOT",
        Path.home() / ".cache/project-application-assistant/legacy-office-conversion",
    )
)
SOFFICE = Path(os.environ.get("SOFFICE_PATH") or shutil.which("soffice") or "soffice")
ENTERPRISE_PATTERN = re.compile(
    r"[\u4e00-\u9fffA-Za-z0-9（）()·—\-\u0020\u3000]{2,80}?"
    r"(?:股份有限公司|有限责任公司|集团有限公司|有限公司|研究院|研究所|学院|中心|"
    r"合作社|事务所|公司|厂)"
    r"(?:[（(][\u4e00-\u9fffA-Za-z0-9·—\-\u0020]{1,20}[）)])?"
    r"(?=$|[\s|,，、;；。])"
)
LIST_ENTITY_LINE_PATTERN = re.compile(
    r"[\u4e00-\u9fffA-Za-z0-9（）()·—\-]{2,100}"
    r"(?:公司|厂|研究院|研究所|学院|中心|合作社|集团|事务所)"
    r"(?:[（(][\u4e00-\u9fffA-Za-z0-9·—\-]{1,20}[）)])?$"
)
HTML_ROW_PATTERN = re.compile(r"<tr\b[^>]*>(.*?)</tr>", re.IGNORECASE | re.DOTALL)
HTML_CELL_PATTERN = re.compile(r"<t[dh]\b[^>]*>(.*?)</t[dh]>", re.IGNORECASE | re.DOTALL)
HTML_TAG_PATTERN = re.compile(r"<[^>]+>")
ENTERPRISE_ALIAS_QUALIFIER_PATTERN = re.compile(
    r"[（(](?:原名|曾用名)\s*[：:]?\s*[^（）()]{1,80}[）)]"
)
YEAR_PATTERN = re.compile(r"(?<!\d)(20\d{2})(?!\d)")
BATCH_PATTERN = re.compile(r"第[一二三四五六七八九十百零〇两\d]+批(?:次)?")
DEPLOYED_PROJECT_INDEX = (
    Path(__file__).resolve().parents[1]
    / "skills"
    / "project-matching"
    / "references"
    / "canonical-project-index.jsonl"
)
SOURCE_PROJECT_INDEX = (
    Path(__file__).resolve().parents[3]
    / "skills"
    / "project-matching"
    / "references"
    / "canonical-project-index.jsonl"
)
DEFAULT_PROJECT_INDEX = Path(
    os.environ.get(
        "JIAOTANG_PROJECT_INDEX_PATH",
        DEPLOYED_PROJECT_INDEX if DEPLOYED_PROJECT_INDEX.is_file() else SOURCE_PROJECT_INDEX,
    )
)
DOCUMENT_STAGE_RULES = (
    ("公示名单", ("公示名单", "拟认定名单", "拟入选名单")),
    ("认定名单", ("认定名单", "通过名单", "入选名单", "企业名单", "项目名单")),
    ("申报通知", ("申报通知", "关于申报", "组织申报", "组织开展", "申报工作")),
    ("管理办法", ("管理办法",)),
    ("实施细则", ("实施细则", "实施办法")),
    ("申报指南", ("申报指南", "工作指南", "工作指引")),
    ("评价标准", ("评价标准", "认定标准", "建设标准")),
    ("征求意见", ("征求意见", "征求意见稿")),
)
PROVINCE_LEVEL_REGIONS = (
    "北京市", "天津市", "上海市", "重庆市",
    "河北省", "山西省", "辽宁省", "吉林省", "黑龙江省", "江苏省", "浙江省",
    "安徽省", "福建省", "江西省", "山东省", "河南省", "湖北省", "湖南省",
    "广东省", "海南省", "四川省", "贵州省", "云南省", "陕西省", "甘肃省",
    "青海省", "台湾省", "内蒙古自治区", "广西壮族自治区", "西藏自治区",
    "宁夏回族自治区", "新疆维吾尔自治区", "香港特别行政区", "澳门特别行政区",
)
ZHEJIANG_CITY_ALIASES = {
    "杭州": "杭州市",
    "宁波": "宁波市",
    "温州": "温州市",
    "嘉兴": "嘉兴市",
    "湖州": "湖州市",
    "绍兴": "绍兴市",
    "金华": "金华市",
    "衢州": "衢州市",
    "舟山": "舟山市",
    "台州": "台州市",
    "丽水": "丽水市",
}
METADATA_RULE_VERSION = "structured-metadata-v2.1.0"
POLICY_CLUSTER_RULE_VERSION = "policy-cluster-v1.0.0"
OFFICIAL_VALIDITY_STAGES = {
    "申报通知",
    "管理办法",
    "实施细则",
    "申报指南",
    "评价标准",
}


def cache_status_reusable(status: str) -> bool:
    return status in {"indexed", "unrecoverable_corrupt"}


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="从知识库清单生成全文索引和云端导入文件")
    parser.add_argument(
        "--manifest",
        type=Path,
        default=Path(os.environ.get("JIAOTANG_MANIFEST_PATH", Path.cwd() / "knowledge-migration/manifest.jsonl")),
    )
    parser.add_argument("--output", type=Path)
    parser.add_argument("--cache", type=Path, default=DEFAULT_CACHE_PATH)
    return parser.parse_args()


def clean_text(text: str) -> str:
    lines = [" ".join(line.split()) for line in text.replace("\x00", " ").splitlines()]
    return "\n".join(line for line in lines if line).strip()[:MAX_TEXT_CHARS]


def normalize_match_text(value: str) -> str:
    normalized = unicodedata.normalize("NFKC", value).lower()
    return re.sub(r"[\s·•,，。；;：:、/\\\-_—()（）\[\]【】《》<>]+", "", normalized)


def metadata_source_name(source: str) -> str:
    """Use only the leaf filename so category directory names cannot contaminate metadata."""
    return Path(source).name


def load_project_catalog(path: Path = DEFAULT_PROJECT_INDEX) -> list[dict[str, object]]:
    if not path.is_file():
        return []
    catalog: list[dict[str, object]] = []
    for line in path.read_text(encoding="utf-8", errors="ignore").splitlines():
        if not line.strip():
            continue
        try:
            record = json.loads(line)
        except json.JSONDecodeError:
            continue
        canonical = str(record.get("canonical_project_name") or "").strip()
        if not canonical:
            continue
        aliases = [canonical, *(str(alias) for alias in record.get("aliases", []))]
        derived_aliases: list[str] = []
        for alias in aliases:
            derived_aliases.extend(
                (
                    re.sub(r"^(?:国家|国家级|省级|市级|区级)", "", alias),
                    re.sub(r"(?:企业认定|企业|项目)$", "", alias),
                    re.sub(
                        r"(?:企业认定|企业|项目)$",
                        "",
                        re.sub(r"^(?:国家|国家级|省级|市级|区级)", "", alias),
                    ),
                )
            )
        aliases.extend(derived_aliases)
        normalized_aliases = sorted(
            {
                normalize_match_text(alias)
                for alias in aliases
                if len(normalize_match_text(alias)) >= 4
            },
            key=len,
            reverse=True,
        )
        catalog.append({**record, "normalized_aliases": normalized_aliases})
    catalog.sort(
        key=lambda item: max((len(alias) for alias in item["normalized_aliases"]), default=0),
        reverse=True,
    )
    return catalog


def infer_canonical_project(
    title: str,
    source: str,
    catalog: list[dict[str, object]],
) -> tuple[str, dict[str, object] | None]:
    haystack = normalize_match_text(f"{title} {metadata_source_name(source)}")
    best: tuple[int, str, dict[str, object], str] | None = None
    for record in catalog:
        for alias in record["normalized_aliases"]:
            if alias not in haystack:
                continue
            candidate = (len(alias), str(record["canonical_project_name"]), record, alias)
            if best is None or candidate[0] > best[0]:
                best = candidate
            break
    if best is None:
        return "", None
    return best[1], {**best[2], "matched_alias": best[3]}


def correct_local_small_giant_scope(
    title: str,
    source: str,
    canonical_project_name: str,
) -> str:
    value = normalize_match_text(f"{title} {metadata_source_name(source)}")
    if any(term in value for term in ("科技小巨人", "创新小巨人", "成长小巨人", "农业科技小巨人")):
        return "地方科技小巨人企业"
    if canonical_project_name != "国家专精特新“小巨人”企业":
        return canonical_project_name
    has_national_anchor = any(term in value for term in ("国家级", "国家专精特新", "工业和信息化部", "工信部"))
    has_local_anchor = any(
        term in value
        for term in (
            "市级专精特新小巨人",
            "省级专精特新小巨人",
            "自治区级专精特新小巨人",
            "省级第三批专精特新小巨人",
            "地方专精特新小巨人",
        )
    )
    has_local_anchor = has_local_anchor or (
        any(term in value for term in ("省级", "市级", "自治区级"))
        and "专精特新" in value
        and "小巨人" in value
    )
    if has_local_anchor and not has_national_anchor:
        return "地方专精特新小巨人企业"
    return canonical_project_name


def infer_policy_year(title: str, source: str) -> int | None:
    for value in (title, metadata_source_name(source)):
        years = [int(year) for year in YEAR_PATTERN.findall(value)]
        if years:
            return max(years)
    return None


def infer_batch(title: str, source: str) -> str:
    match = BATCH_PATTERN.search(f"{title} {metadata_source_name(source)}")
    return match.group(0) if match else ""


NATIONAL_SMALL_GIANT_BATCH_YEARS = {
    "第一批": 2019,
    "第二批": 2020,
    "第三批": 2021,
    "第四批": 2022,
    "第五批": 2023,
    "第六批": 2024,
    "第七批": 2025,
    "第八批": 2026,
}


def infer_small_giant_batch_year(
    title: str,
    source: str,
    canonical_project_name: str,
    batch: str,
) -> int | None:
    value = f"{title} {metadata_source_name(source)}"
    if canonical_project_name != "国家专精特新“小巨人”企业":
        return None
    if "复核" in value or "重点" in value:
        return None
    if "专精特新" not in value or "小巨人" not in value or "名单" not in value:
        return None
    return NATIONAL_SMALL_GIANT_BATCH_YEARS.get(batch)


def infer_single_province_list_region(
    content: str,
    canonical_project_name: str,
    document_role: str,
) -> str:
    if canonical_project_name != "国家专精特新“小巨人”企业":
        return ""
    if document_role != "50_名单与对标":
        return ""
    regions = [region for region in PROVINCE_LEVEL_REGIONS if region in content]
    return regions[0] if len(regions) == 1 else ""


def infer_document_stage(title: str, source: str, document_role: str) -> str:
    value = f"{title} {metadata_source_name(source)}"
    if "征求意见" in value:
        return "征求意见"
    if "名单" in value and any(term in value for term in ("公示", "拟认定", "拟入选")):
        return "公示名单"
    if "名单" in value and any(term in value for term in ("认定", "通过", "入选")):
        return "认定名单"
    if "公示稿" in value or "公示" in value:
        return "公示稿" if "公示稿" in value else "公示"
    for stage, terms in DOCUMENT_STAGE_RULES:
        if any(term in value for term in terms):
            return stage
    if document_role == "50_名单与对标" and "名单" in value:
        return "名单"
    if "通知" in value:
        return "通知"
    return "其他"


def infer_region(
    title: str,
    source: str,
    project_record: dict[str, object] | None,
    catalog: list[dict[str, object]],
) -> str:
    value = f"{title} {metadata_source_name(source)}"
    known_regions = {
        str(region)
        for record in catalog
        for region in record.get("regions", [])
        if str(region) and str(region) != "待确认"
    }
    known_regions.update({"全国", *PROVINCE_LEVEL_REGIONS, *ZHEJIANG_CITY_ALIASES.values()})
    matches = sorted(
        {region for region in known_regions if region in value},
        key=lambda region: (-len(region), region),
    )
    for alias, canonical in ZHEJIANG_CITY_ALIASES.items():
        if alias in value and canonical not in matches:
            matches.append(canonical)
    if not matches and project_record:
        matches = [
            str(region)
            for region in project_record.get("regions", [])
            if str(region) and str(region) != "待确认"
        ]
    return "|".join(dict.fromkeys(matches))


def infer_validity_status(title: str, content: str) -> str:
    try:
        from scripts.build_policy_version_links import detect_policy_status, extract_lifecycle_evidence
    except ModuleNotFoundError:
        from build_policy_version_links import detect_policy_status, extract_lifecycle_evidence

    lifecycle = extract_lifecycle_evidence(content[:12_000])
    if lifecycle["self_invalid"]:
        return "invalid"
    draft_text = f"{title}\n{content[:12_000]}"
    if any(term in draft_text for term in ("公示稿", "征求意见稿", "草案")):
        return "draft"
    return detect_policy_status(title)


def infer_policy_replacement(
    title: str,
    source: str,
    content: str,
    document_role: str,
) -> dict[str, str]:
    value = f"{title}\n{metadata_source_name(source)}\n{content[:20_000]}"
    result = {
        "validity_status": infer_validity_status(title, content),
        "replacement_title": "",
        "replacement_basis": "",
        "replacement_url": "",
    }
    explicit_historical_path = any(
        term in source for term in ("历史政策", "历史规则", "历史培训")
    )
    if result["validity_status"] == "draft" and not explicit_historical_path:
        return result
    title_and_source = f"{title}\n{metadata_source_name(source)}"
    if (
        "杭市管〔2025〕93号" in value
        or "杭州市知识产权强企认定管理办法" in title_and_source
    ):
        result.update(
            {
                "validity_status": "active_candidate",
                "replacement_title": "",
                "replacement_basis": "",
                "replacement_url": "",
            }
        )
        return result
    old_hangzhou_patent_program = (
        "杭市管〔2020〕38号" in value
        or "杭州市专利试点企业和示范企业认定管理办法" in title_and_source
        or (
            any(term in title_and_source for term in ("专利试点", "专利示范"))
            and "知识产权强企" not in title_and_source
        )
    )
    if old_hangzhou_patent_program:
        result.update(
            {
                "validity_status": (
                    "superseded"
                    if any(
                        term in title_and_source
                        for term in ("管理办法", "杭市管〔2020〕38号")
                    )
                    else "historical_reference"
                ),
                "replacement_title": "《杭州市知识产权强企认定管理办法》（杭市管〔2025〕93号）",
                "replacement_basis": "杭市管〔2025〕93号第十四条明确自2025年9月15日起施行，杭市管〔2020〕38号同时废止；旧表单、清单和培训资料仅作历史参考。",
                "replacement_url": "https://scjg.hangzhou.gov.cn/art/2025/8/19/art_1229144701_1858746.html",
            }
        )
        return result
    if "杭州市AI工厂" in title_and_source:
        result.update(
            {
                "validity_status": "active_candidate",
                "replacement_title": "",
                "replacement_basis": "",
                "replacement_url": "",
            }
        )
        return result
    old_hangzhou_future_factory = (
        "未来工厂" in title_and_source
        and (
            "杭州市" in title_and_source
            or "杭州市未来工厂" in source
        )
        and "浙江省未来工厂" not in title_and_source
    )
    if old_hangzhou_future_factory:
        result.update(
            {
                "validity_status": "historical_reference",
                "replacement_title": "杭州市AI工厂",
                "replacement_basis": "杭州市2026年市级培育申报入口已转为AI工厂；旧杭州市未来工厂通知、名单、申报书和培训资料保留为历史参考。浙江省未来工厂仍为独立省级体系，不受本规则影响。",
                "replacement_url": "https://zfgb.hangzhou.gov.cn/148/102220263/t103220263024/530188.shtml",
            }
        )
        return result
    if explicit_historical_path:
        result.update(
            {
                "validity_status": "historical_reference",
                "replacement_title": result["replacement_title"],
                "replacement_basis": (
                    result["replacement_basis"]
                    or "文件已进入明确的历史资料层，只能用于追溯历史规则、表单或方法，不得作为当前申报依据。"
                ),
                "replacement_url": result["replacement_url"],
            }
        )
        return result
    if "工信部企业〔2026〕2号" in value or (
        "优质中小企业梯度培育管理办法" in value
        and "科技和创新型中小企业" in value
        and "质量评价得分" in value
    ):
        result["validity_status"] = "active_candidate"
        return result
    if "工信部企业〔2022〕63号" in value or "优质中小企业梯度培育管理暂行办法" in value:
        result.update(
            {
                "validity_status": "superseded",
                "replacement_title": "《优质中小企业梯度培育管理办法》（工信部企业〔2026〕2号）",
                "replacement_basis": "2026年新办法已发布，旧办法只保留为历史档案。2026年度小巨人复核曾按工信厅企业函〔2026〕117号使用旧标准，但该批复核已经结束；旧办法不得用于当前或未来的新申报、复核、评分和材料写作，也不得补充现行标准没有规定的条件。",
                "replacement_url": "https://www.miit.gov.cn/zwgk/zcwj/wjfb/tz/art/2026/art_5546b1a622ea4d73bd2ca395b73bd4eb.html",
            }
        )
        return result
    if (
        any(term in value for term in ("专精特新", "小巨人"))
        and document_role in {"20_申报指南与规则", "40_内部培训与方法", "60_申报案例与建设方案"}
        and "2026" not in f"{title} {metadata_source_name(source)}"
    ):
        result.update(
            {
                "validity_status": "historical_reference",
                "replacement_title": "《优质中小企业梯度培育管理办法》（工信部企业〔2026〕2号）",
                "replacement_basis": "旧指南、培训材料和申请案例仅供历史结构参考，不得作为现行认定条件。",
                "replacement_url": "https://www.miit.gov.cn/zwgk/zcwj/wjfb/tz/art/2026/art_5546b1a622ea4d73bd2ca395b73bd4eb.html",
            }
        )
    if "浙江省高新技术企业研究开发中心" in value or "浙江省研发中心.ppt" in value:
        result.update(
            {
                "validity_status": "superseded",
                "replacement_title": "浙江省企业研究院",
                "replacement_basis": "浙经信高新〔2025〕169号建立省重点企业研究院、省企业研究院新体系；原省高新技术企业研究开发中心纳入省企业研究院序列，不再重复申报认定。",
                "replacement_url": "https://zj87.jxt.zj.gov.cn/zlzq/web/views/article/news/detail.html?id=280535",
            }
        )
    if "杭州市研发中心" in title or "市研发政策解读" in title:
        result.update(
            {
                "validity_status": "historical_reference",
                "replacement_title": "杭州市企业高新技术研究开发中心",
                "replacement_basis": "内部培训简称与旧讲义不得替代杭科高〔2022〕39号正式管理办法；该正式办法在杭州市科技局2025年清理结果中列为继续有效。",
                "replacement_url": "https://zfgb.hangzhou.gov.cn/11/105220253/t117220253054/518912.shtml",
            }
        )
    if "杭科高〔2022〕39号" in value or "杭州市企业高新技术研究开发中心管理办法" in title:
        result.update(
            {
                "validity_status": "active_candidate",
                "replacement_title": "",
                "replacement_basis": "",
                "replacement_url": "",
            }
        )
    return result


def _confirmed_alias_correction(
    title: str,
    source: str,
    metadata: dict[str, object],
    corrections: list[dict[str, object]],
) -> dict[str, object] | None:
    haystack = normalize_match_text(f"{title} {metadata_source_name(source)}")
    policy_year = metadata.get("policy_year")
    regions = set(str(metadata.get("region") or "").split("|"))
    candidates: list[tuple[int, dict[str, object]]] = []
    for correction in corrections:
        if str(correction.get("status") or "") != "confirmed":
            continue
        raw_name = str(correction.get("raw_project_name") or "").strip()
        canonical_name = str(correction.get("canonical_project_name") or "").strip()
        normalized_raw = normalize_match_text(raw_name)
        if not normalized_raw or not canonical_name or normalized_raw not in haystack:
            continue
        correction_region = str(correction.get("region") or "").strip()
        if correction_region and correction_region not in regions:
            continue
        start_year = correction.get("start_year")
        end_year = correction.get("end_year")
        if (start_year or end_year) and not policy_year:
            continue
        if policy_year and start_year and int(policy_year) < int(start_year):
            continue
        if policy_year and end_year and int(policy_year) > int(end_year):
            continue
        candidates.append((len(normalized_raw), correction))
    return max(candidates, key=lambda item: item[0])[1] if candidates else None


def _metadata_evidence(
    title: str,
    source: str,
    content: str,
    metadata: dict[str, object],
    correction: dict[str, object] | None,
    document_role: str,
) -> list[dict[str, object]]:
    title_filename = f"{title} | {metadata_source_name(source)}"[:800]
    evidence: list[dict[str, object]] = []
    field_settings = {
        "canonical_project_name": (
            "manual_alias" if correction else "catalog_alias",
            str(correction.get("raw_project_name") or "") if correction else str(metadata.get("project_matched_term") or ""),
            "confirmed" if correction else ("high" if metadata.get("canonical_project_name") else "low"),
        ),
        "region": (
            "explicit_region_or_catalog_scope",
            str(metadata.get("region") or ""),
            "high" if metadata.get("region") and str(metadata.get("region")) in title_filename else ("medium" if metadata.get("region") else "low"),
        ),
        "document_stage": (
            "stage_term",
            str(metadata.get("document_stage") or ""),
            "high" if metadata.get("document_stage") != "其他" else "low",
        ),
        "validity_status": (
            "policy_replacement_rule" if metadata.get("replacement_url") else "lifecycle_term",
            str(metadata.get("replacement_title") or metadata.get("validity_status") or ""),
            "high" if metadata.get("validity_status") in {"invalid", "superseded"} else "medium",
        ),
        "policy_year": (
            "year_pattern",
            str(metadata.get("policy_year") or ""),
            "high" if metadata.get("policy_year") else "low",
        ),
        "batch": (
            "batch_pattern",
            str(metadata.get("batch") or ""),
            "high" if metadata.get("batch") else "low",
        ),
    }
    structured_role = document_role.startswith(("10_", "20_", "50_"))
    if not structured_role:
        for field_name in ("document_stage", "policy_year", "batch"):
            field_settings.pop(field_name)
    elif not document_role.startswith("50_"):
        field_settings.pop("batch")
    for field_name, (method, matched_term, confidence) in field_settings.items():
        source_excerpt = title_filename
        if field_name == "validity_status" and content:
            source_excerpt = clean_text(content[:800]) or title_filename
        if field_name == "validity_status" and metadata.get("replacement_url"):
            source_excerpt = " | ".join(
                (
                    str(metadata.get("replacement_basis") or ""),
                    str(metadata.get("replacement_url") or ""),
                )
            )[:800]
        evidence.append(
            {
                "field_name": field_name,
                "inferred_value": str(metadata.get(field_name) or ""),
                "matched_term": matched_term,
                "match_method": method,
                "source_scope": "manual_correction" if correction and field_name == "canonical_project_name" else ("official_rule" if field_name == "validity_status" and metadata.get("replacement_url") else ("content" if field_name == "validity_status" else "title_filename")),
                "source_excerpt": source_excerpt,
                "rule_version": METADATA_RULE_VERSION,
                "confidence": confidence,
                "review_status": "confirmed" if correction and field_name == "canonical_project_name" else ("needs_review" if confidence == "low" else "unreviewed"),
                "correction_id": correction.get("id") if correction else None,
            }
        )
    return evidence


def infer_document_metadata(
    title: str,
    source: str,
    content: str,
    document_role: str,
    catalog: list[dict[str, object]] | None = None,
    corrections: list[dict[str, object]] | None = None,
) -> dict[str, object]:
    project_catalog = catalog if catalog is not None else load_project_catalog()
    canonical_project_name, project_record = infer_canonical_project(
        title, source, project_catalog
    )
    canonical_project_name = correct_local_small_giant_scope(
        title, source, canonical_project_name
    )
    replacement = infer_policy_replacement(title, source, content, document_role)
    batch = infer_batch(title, source)
    policy_year = infer_policy_year(title, source)
    if policy_year is None:
        policy_year = infer_small_giant_batch_year(
            title,
            source,
            canonical_project_name,
            batch,
        )
    region = infer_region(title, source, project_record, project_catalog)
    list_region = infer_single_province_list_region(
        content,
        canonical_project_name,
        document_role,
    )
    if list_region:
        region = list_region
    project_id = infer_project_id(canonical_project_name, project_record)
    case_pack = infer_case_pack_metadata(source, document_role, project_id, policy_year, batch)
    metadata: dict[str, object] = {
        "project_id": project_id,
        **case_pack,
        "canonical_project_name": canonical_project_name,
        "region": region,
        "document_stage": infer_document_stage(title, source, document_role),
        "validity_status": replacement["validity_status"],
        "policy_year": policy_year,
        "batch": batch,
        "replacement_title": replacement["replacement_title"],
        "replacement_basis": replacement["replacement_basis"],
        "replacement_url": replacement["replacement_url"],
        "project_matched_term": str(project_record.get("matched_alias") or "") if project_record else "",
    }
    correction = _confirmed_alias_correction(
        title, source, metadata, corrections or []
    )
    if correction:
        metadata["canonical_project_name"] = str(
            correction["canonical_project_name"]
        )
    metadata["match_evidence"] = _metadata_evidence(
        title, source, content, metadata, correction, document_role
    )
    return metadata


KNOWN_PROJECT_IDS = {
    "杭州市企业技术中心": "hangzhou-enterprise-technology-center",
    "浙江省企业技术中心": "zhejiang-enterprise-technology-center",
    "浙江省重点企业研究院": "zhejiang-enterprise-institute",
    "浙江省企业研究院": "zhejiang-enterprise-institute",
}


def infer_project_id(
    canonical_project_name: str, project_record: dict[str, object] | None
) -> str:
    if canonical_project_name in KNOWN_PROJECT_IDS:
        return KNOWN_PROJECT_IDS[canonical_project_name]
    if not canonical_project_name:
        return ""
    explicit = str((project_record or {}).get("canonical_project_id") or "").strip()
    if explicit:
        return explicit
    digest = hashlib.sha256(canonical_project_name.encode("utf-8")).hexdigest()[:16]
    return f"project-{digest}"


def infer_case_pack_metadata(
    source: str,
    document_role: str,
    project_id: str,
    policy_year: int | None,
    batch: str,
) -> dict[str, object]:
    path = Path(source)
    parts = path.parts
    try:
        case_index = next(
            index for index, value in enumerate(parts) if value.startswith("60_申报案例与建设方案")
        )
    except StopIteration:
        case_index = -1
    in_case_library = document_role.startswith("60_") or case_index >= 0
    if not in_case_library:
        return {
            "case_pack_id": "",
            "document_type": infer_case_document_type(path.name, document_role),
            "evidence_type": "",
            "upload_action": "exclude" if document_role.startswith("90_") else "review",
            "verification_status": "unreviewed",
            "case_pack_title": "",
            "case_pack_source_root": "",
        }
    relative = parts[case_index + 1 :] if case_index >= 0 else parts[-3:]
    group_parts = relative[:-1]
    group = "/".join(group_parts[:3]) or path.stem
    identity = "|".join((project_id, str(policy_year or ""), batch, group))
    case_pack_id = f"case-{hashlib.sha256(identity.encode('utf-8')).hexdigest()[:20]}"
    document_type = infer_case_document_type(path.name, document_role)
    evidence_type = infer_case_evidence_type(path.name, document_type)
    return {
        "case_pack_id": case_pack_id,
        "document_type": document_type,
        "evidence_type": evidence_type,
        "upload_action": "local_only" if document_role.startswith("90_") else "review",
        "verification_status": "auto_grouped",
        "case_pack_title": group_parts[-1] if group_parts else path.stem,
        "case_pack_source_root": "/".join(group_parts),
    }


def infer_case_document_type(file_name: str, document_role: str) -> str:
    value = f"{file_name} {document_role}"
    for terms, result in (
        (("申报书", "申请书"), "application"),
        (("建设方案",), "construction_plan"),
        (("可研", "技术方案"), "technical_plan"),
        (("专家意见", "评审意见"), "expert_review"),
        (("公示", "认定", "结果"), "recognition_result"),
        (("附件", "佐证", "证明"), "evidence_attachment"),
    ):
        if any(term in value for term in terms):
            return result
    return "reference_document"


def infer_case_evidence_type(file_name: str, document_type: str) -> str:
    if document_type != "evidence_attachment":
        return ""
    for terms, result in (
        (("财务", "审计", "纳税"), "financial"),
        (("研发人员", "人员名册"), "personnel"),
        (("设备", "场地"), "equipment_site"),
        (("专利", "软著", "知识产权"), "intellectual_property"),
        (("制度",), "management_system"),
    ):
        if any(term in file_name for term in terms):
            return result
    return "other_evidence"


def normalize_policy_document_number(value: str) -> str:
    match = re.search(
        r"([A-Za-z\u4e00-\u9fff]{1,12})\s*[\u3014\[\u3010（(]\s*((?:19|20)\d{2})\s*[\u3015\]\u3011）)]\s*(\d{1,5})\s*号",
        value,
    )
    if not match:
        return ""
    authority = re.sub(r"[^A-Za-z\u4e00-\u9fff]", "", match.group(1))
    return f"{authority}〔{match.group(2)}〕{int(match.group(3))}号"


def normalize_policy_cluster_title(value: str) -> str:
    title = Path(value).stem
    title = re.sub(r"^(?:附件|[附件一二三四五六七八九十\d]+)[：:]\s*", "", title)
    title = re.sub(r"(?:副本|复制件|扫描件|转换版)$", "", title)
    return re.sub(r"[^0-9A-Za-z\u4e00-\u9fff]+", "", title).lower()


def generic_policy_cluster_title(normalized_title: str) -> bool:
    generic_titles = {
        "目录", "附件", "申报材料", "申报要求", "申报办法",
        "人员网上申报办法", "申报人员网上申报办法", "申报指南",
        "管理办法", "实施细则", "通知", "公示", "名单",
    }
    return normalized_title in generic_titles


def ensure_policy_cluster_schema(connection: sqlite3.Connection) -> None:
    connection.executescript(
        """
        CREATE TABLE IF NOT EXISTS policy_document_clusters (
            id INTEGER PRIMARY KEY,
            cluster_key TEXT NOT NULL UNIQUE,
            normalized_title TEXT NOT NULL,
            document_number TEXT NOT NULL DEFAULT '',
            canonical_project_name TEXT NOT NULL DEFAULT '',
            region TEXT NOT NULL DEFAULT '',
            policy_year INTEGER,
            representative_document_id INTEGER NOT NULL REFERENCES documents(id),
            match_method TEXT NOT NULL,
            confidence TEXT NOT NULL,
            rule_version TEXT NOT NULL,
            created_at TEXT NOT NULL,
            updated_at TEXT NOT NULL
        );
        CREATE INDEX IF NOT EXISTS policy_document_clusters_number_idx
            ON policy_document_clusters(document_number,policy_year,region);
        CREATE TABLE IF NOT EXISTS policy_document_cluster_members (
            id INTEGER PRIMARY KEY,
            cluster_id INTEGER NOT NULL REFERENCES policy_document_clusters(id),
            document_id INTEGER NOT NULL UNIQUE REFERENCES documents(id),
            membership_basis TEXT NOT NULL,
            confidence TEXT NOT NULL,
            created_at TEXT NOT NULL
        );
        CREATE INDEX IF NOT EXISTS policy_document_cluster_members_cluster_idx
            ON policy_document_cluster_members(cluster_id,document_id);
        CREATE TABLE IF NOT EXISTS policy_cluster_manual_assignments (
            document_id INTEGER PRIMARY KEY REFERENCES documents(id),
            manual_cluster_key TEXT NOT NULL,
            operation_type TEXT NOT NULL,
            note TEXT NOT NULL DEFAULT '',
            updated_by TEXT NOT NULL,
            updated_at TEXT NOT NULL
        );
        CREATE INDEX IF NOT EXISTS policy_cluster_manual_assignments_key_idx
            ON policy_cluster_manual_assignments(manual_cluster_key,document_id);
        CREATE TABLE IF NOT EXISTS policy_cluster_manual_operations (
            id INTEGER PRIMARY KEY,
            operation_type TEXT NOT NULL,
            source_cluster_ids TEXT NOT NULL,
            document_ids TEXT NOT NULL,
            target_manual_cluster_key TEXT NOT NULL,
            note TEXT NOT NULL DEFAULT '',
            operated_by TEXT NOT NULL,
            operated_at TEXT NOT NULL,
            previous_assignments TEXT NOT NULL DEFAULT '{}',
            undone_by TEXT NOT NULL DEFAULT '',
            undone_at TEXT
        );
        CREATE INDEX IF NOT EXISTS policy_cluster_manual_operations_time_idx
            ON policy_cluster_manual_operations(operated_at,id);
        CREATE TABLE IF NOT EXISTS policy_verification_propagations (
            id INTEGER PRIMARY KEY,
            source_queue_id INTEGER NOT NULL REFERENCES policy_verification_queue(id),
            cluster_id INTEGER NOT NULL REFERENCES policy_document_clusters(id),
            source_document_id INTEGER NOT NULL REFERENCES documents(id),
            target_document_id INTEGER NOT NULL REFERENCES documents(id),
            field_name TEXT NOT NULL,
            propagated_value TEXT NOT NULL,
            official_source_url TEXT NOT NULL DEFAULT '',
            evidence_excerpt TEXT NOT NULL,
            rule_version TEXT NOT NULL,
            propagated_by TEXT NOT NULL,
            propagated_at TEXT NOT NULL,
            UNIQUE(source_queue_id,target_document_id,field_name)
        );
        CREATE INDEX IF NOT EXISTS policy_verification_propagations_target_idx
            ON policy_verification_propagations(target_document_id,propagated_at);
        """
    )
    operation_columns = {
        row[1] for row in connection.execute("PRAGMA table_info(policy_cluster_manual_operations)")
    }
    if "previous_assignments" not in operation_columns:
        connection.execute(
            "ALTER TABLE policy_cluster_manual_operations ADD COLUMN previous_assignments TEXT NOT NULL DEFAULT '{}'"
        )
    if "undone_by" not in operation_columns:
        connection.execute(
            "ALTER TABLE policy_cluster_manual_operations ADD COLUMN undone_by TEXT NOT NULL DEFAULT ''"
        )
    if "undone_at" not in operation_columns:
        connection.execute(
            "ALTER TABLE policy_cluster_manual_operations ADD COLUMN undone_at TEXT"
        )


def rebuild_policy_document_clusters(connection: sqlite3.Connection) -> dict[str, int]:
    ensure_policy_cluster_schema(connection)
    connection.row_factory = sqlite3.Row
    connection.execute(
        """
        DELETE FROM policy_cluster_manual_assignments
        WHERE document_id NOT IN (SELECT id FROM documents)
        """
    )
    manual_assignments = {
        int(row["document_id"]): dict(row)
        for row in connection.execute(
            "SELECT * FROM policy_cluster_manual_assignments"
        ).fetchall()
    }
    rows = connection.execute(
        """
        SELECT id,title,source,document_role,canonical_project_name,
               region,policy_year
        FROM documents
        WHERE document_role LIKE '10_%' OR document_role LIKE '20_%'
        ORDER BY id
        """
    ).fetchall()
    grouped: dict[str, list[dict[str, object]]] = {}
    for row in rows:
        item = dict(row)
        title = str(item.get("title") or "")
        source_name = Path(str(item.get("source") or "")).name
        document_number = normalize_policy_document_number(f"{title}\n{source_name}")
        normalized_title = normalize_policy_cluster_title(title)
        if document_number:
            identity = f"number:{document_number}"
            match_method, confidence = "document_number", "high"
        elif len(normalized_title) >= 8 and not generic_policy_cluster_title(normalized_title):
            scope = "|".join(
                (
                    normalized_title,
                    str(item.get("canonical_project_name") or ""),
                    str(item.get("region") or ""),
                    str(item.get("policy_year") or ""),
                )
            )
            identity = f"title:{hashlib.sha256(scope.encode('utf-8')).hexdigest()}"
            match_method, confidence = "exact_normalized_title_scope", "medium"
        else:
            identity = f"document:{int(item['id'])}"
            match_method, confidence = "singleton", "low"
        manual_assignment = manual_assignments.get(int(item["id"]))
        if manual_assignment:
            identity = str(manual_assignment["manual_cluster_key"])
            match_method = f"manual_{manual_assignment['operation_type']}"
            confidence = "high"
        item.update(
            {
                "cluster_key": identity,
                "normalized_title": normalized_title,
                "document_number": document_number,
                "match_method": match_method,
                "cluster_confidence": confidence,
            }
        )
        grouped.setdefault(identity, []).append(item)

    now = datetime.now(timezone.utc).isoformat()
    active_cluster_ids: list[int] = []
    connection.execute("DELETE FROM policy_document_cluster_members")
    for cluster_key, members in grouped.items():
        representative = members[0]
        connection.execute(
            """
            INSERT INTO policy_document_clusters(
                cluster_key,normalized_title,document_number,canonical_project_name,
                region,policy_year,representative_document_id,match_method,confidence,
                rule_version,created_at,updated_at
            ) VALUES (?,?,?,?,?,?,?,?,?,?,?,?)
            ON CONFLICT(cluster_key) DO UPDATE SET
                normalized_title=excluded.normalized_title,
                document_number=excluded.document_number,
                canonical_project_name=excluded.canonical_project_name,
                region=excluded.region,
                policy_year=excluded.policy_year,
                representative_document_id=excluded.representative_document_id,
                match_method=excluded.match_method,
                confidence=excluded.confidence,
                rule_version=excluded.rule_version,
                updated_at=excluded.updated_at
            """,
            (
                cluster_key,
                representative["normalized_title"],
                representative["document_number"],
                str(representative.get("canonical_project_name") or ""),
                str(representative.get("region") or ""),
                representative.get("policy_year"),
                int(representative["id"]),
                representative["match_method"],
                representative["cluster_confidence"],
                POLICY_CLUSTER_RULE_VERSION,
                now,
                now,
            ),
        )
        cluster_id = int(
            connection.execute(
                "SELECT id FROM policy_document_clusters WHERE cluster_key=?", (cluster_key,)
            ).fetchone()[0]
        )
        active_cluster_ids.append(cluster_id)
        connection.executemany(
            """
            INSERT INTO policy_document_cluster_members(
                cluster_id,document_id,membership_basis,confidence,created_at
            ) VALUES (?,?,?,?,?)
            """,
            (
                (
                    cluster_id,
                    int(member["id"]),
                    str(member["match_method"]),
                    str(member["cluster_confidence"]),
                    now,
                )
                for member in members
            ),
        )
    if active_cluster_ids:
        placeholders = ",".join("?" for _ in active_cluster_ids)
        connection.execute(
            f"""
            DELETE FROM policy_document_clusters
            WHERE id NOT IN ({placeholders})
              AND id NOT IN (SELECT DISTINCT cluster_id FROM policy_verification_propagations)
            """,
            active_cluster_ids,
        )
    duplicate_clusters = sum(1 for members in grouped.values() if len(members) > 1)
    duplicate_documents = sum(len(members) for members in grouped.values() if len(members) > 1)
    return {
        "clusters": len(grouped),
        "duplicate_clusters": duplicate_clusters,
        "duplicate_documents": duplicate_documents,
    }


def insert_metadata_audit_records(
    connection: sqlite3.Connection,
    document_id: int,
    document_role: str,
    metadata: dict[str, object],
) -> None:
    for item in metadata.get("match_evidence", []):
        connection.execute(
            """
            INSERT OR REPLACE INTO metadata_match_evidence(
                document_id,field_name,inferred_value,matched_term,match_method,
                source_scope,source_excerpt,rule_version,confidence,review_status,
                correction_id,created_at
            ) VALUES (?,?,?,?,?,?,?,?,?,?,?,?)
            """,
            (
                document_id,
                item["field_name"],
                item["inferred_value"],
                item["matched_term"],
                item["match_method"],
                item["source_scope"],
                item["source_excerpt"],
                item["rule_version"],
                item["confidence"],
                item["review_status"],
                item["correction_id"],
                datetime.now(timezone.utc).isoformat(),
            ),
        )
    if not document_role.startswith(("10_", "20_")):
        return
    queue_items: list[tuple[str, str]] = []
    if not metadata.get("canonical_project_name") and metadata.get("document_stage") != "其他":
        queue_items.append(("项目名称未能可靠映射", "medium"))
    if (
        metadata.get("document_stage") in OFFICIAL_VALIDITY_STAGES
        and metadata.get("validity_status")
        in {"active_candidate", "revised", "trial", "draft"}
    ):
        queue_items.append(("有效性需要官方网站复核", "high"))
    for reason, priority in queue_items:
        connection.execute(
            """
            INSERT OR IGNORE INTO policy_verification_queue(
                document_id,reason,priority,status,created_at,updated_at
            ) VALUES (?,?,?,'pending',?,?)
            """,
            (
                document_id,
                reason,
                priority,
                datetime.now(timezone.utc).isoformat(),
                datetime.now(timezone.utc).isoformat(),
            ),
        )


def extract_pdf(path: Path) -> str:
    parts: list[str] = []
    character_count = 0
    try:
        import fitz
    except ImportError:
        from pypdf import PdfReader

        document = PdfReader(path)
        for page in document.pages:
            page_text = page.extract_text() or ""
            parts.append(page_text)
            character_count += len(page_text)
            if character_count >= MAX_TEXT_CHARS:
                break
    else:
        fitz.TOOLS.mupdf_display_errors(False)
        fitz.TOOLS.mupdf_display_warnings(False)
        with fitz.open(path) as document:
            for page in document:
                page_text = page.get_text("text")
                parts.append(page_text)
                character_count += len(page_text)
                if character_count >= MAX_TEXT_CHARS:
                    break
    return clean_text("\n".join(parts))


def extract_pdf_isolated(path: Path) -> str:
    process = subprocess.run(
        [sys.executable, str(Path(__file__).resolve()), "--extract-pdf-worker", str(path)],
        capture_output=True,
        check=False,
        timeout=45,
    )
    if process.returncode != 0:
        return ""
    return process.stdout.decode("utf-8", errors="ignore")[:MAX_TEXT_CHARS]


def extract_docx(path: Path) -> str:
    from docx import Document

    document = Document(path)
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return clean_text("\n".join(parts))


def extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(path.open("rb"), data_only=True, read_only=True)
    parts: list[str] = []
    character_count = 0
    try:
        for sheet in workbook.worksheets:
            parts.append(f"工作表：{sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                line = " | ".join("" if value is None else str(value) for value in row)
                if line.strip(" |"):
                    parts.append(line)
                    character_count += len(line)
                if character_count >= MAX_TEXT_CHARS:
                    return clean_text("\n".join(parts))
    finally:
        workbook.close()
    return clean_text("\n".join(parts))


def extract_xls(path: Path) -> str:
    import xlrd

    workbook = xlrd.open_workbook(path, on_demand=True)
    parts: list[str] = []
    character_count = 0
    try:
        for sheet_name in workbook.sheet_names():
            sheet = workbook.sheet_by_name(sheet_name)
            parts.append(f"工作表：{sheet_name}")
            for row_index in range(sheet.nrows):
                line = " | ".join(str(value) for value in sheet.row_values(row_index))
                parts.append(line)
                character_count += len(line)
                if character_count >= MAX_TEXT_CHARS:
                    return clean_text("\n".join(parts))
    finally:
        workbook.release_resources()
    return clean_text("\n".join(parts))


def extract_pptx(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(path)
    parts: list[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text)
        if slide_parts:
            parts.append(f"第{slide_number}页\n" + "\n".join(slide_parts))
    return clean_text("\n".join(parts))


def extract_legacy_office(path: Path) -> str:
    process = subprocess.run(
        ["/usr/bin/textutil", "-convert", "txt", "-stdout", str(path)],
        capture_output=True,
        check=False,
        timeout=120,
    )
    text = clean_text(process.stdout.decode("utf-8", errors="ignore"))
    if process.returncode == 0 and len(text) >= MIN_TEXT_CHARS:
        return text

    digest = hashlib.sha256(str(path).encode("utf-8")).hexdigest()
    output = LEGACY_CONVERSION_ROOT / digest
    output.mkdir(parents=True, exist_ok=True)
    target_extension = ".pptx" if path.suffix.lower() == ".ppt" else ".docx"
    converted = output / f"{path.stem}{target_extension}"
    if not converted.exists() and SOFFICE.exists():
        subprocess.run(
            [
                str(SOFFICE), "--headless", "--convert-to", target_extension.lstrip("."),
                "--outdir", str(output), str(path),
            ],
            capture_output=True,
            check=False,
            timeout=240,
        )
    if not converted.exists():
        return ""
    if target_extension == ".pptx":
        return extract_pptx(converted)
    return extract_docx(converted)


def extract_plain(path: Path) -> str:
    return clean_text(path.read_text(encoding="utf-8", errors="ignore"))


def flatten_plist(value: object) -> list[str]:
    if isinstance(value, dict):
        parts: list[str] = []
        for key, child in value.items():
            parts.append(str(key))
            parts.extend(flatten_plist(child))
        return parts
    if isinstance(value, list):
        parts = []
        for child in value:
            parts.extend(flatten_plist(child))
        return parts
    if isinstance(value, bytes):
        return [value.decode("utf-8", errors="ignore")]
    return [str(value)]


def extract_manual(path: Path) -> tuple[str, str]:
    if path.name == ".WeDrive":
        return "", "non_content_placeholder"
    if ".obsidian/plugins" in path.as_posix():
        return "", "non_content_placeholder"
    extension = path.suffix.lower()
    if extension in {".js", ".css"}:
        text = extract_plain(path)
        return text, "indexed" if len(text) >= MIN_TEXT_CHARS else "empty_non_content"
    if extension == ".textclipping":
        try:
            text = clean_text("\n".join(flatten_plist(plistlib.loads(path.read_bytes()))))
        except Exception:
            text = ""
        return text, "indexed" if len(text) >= MIN_TEXT_CHARS else "empty_non_content"
    if extension == ".emmx":
        parts: list[str] = []
        try:
            with zipfile.ZipFile(path) as archive:
                for name in archive.namelist():
                    if name.lower().endswith((".xml", ".json", ".txt")):
                        parts.append(archive.read(name).decode("utf-8", errors="ignore"))
        except zipfile.BadZipFile:
            pass
        text = clean_text("\n".join(parts))
        return text, "indexed" if len(text) >= MIN_TEXT_CHARS else "empty_non_content"

    header = path.read_bytes()[:8]
    if header.startswith(b"%PDF"):
        text = extract_pdf_isolated(path)
        return text, "indexed" if len(text) >= MIN_TEXT_CHARS else "ocr_required"
    if header.startswith(b"PK"):
        try:
            text = extract_docx(path)
        except Exception:
            text = ""
        return text, "indexed" if len(text) >= MIN_TEXT_CHARS else "convert_required"
    if header.startswith(bytes.fromhex("d0cf11e0")):
        text = extract_legacy_office(path)
        return text, "indexed" if len(text) >= MIN_TEXT_CHARS else "convert_required"
    return "", "non_content_manual_review"


def extract(path: Path, extension: str) -> tuple[str, str]:
    header = path.read_bytes()[:8]
    if extension == ".pdf":
        text = extract_pdf_isolated(path)
        return text, "indexed" if len(text) >= MIN_TEXT_CHARS else "ocr_required"
    if extension in {".docx", ".docm"}:
        if header.startswith(bytes.fromhex("d0cf11e0")) or extension == ".docm":
            text = extract_legacy_office(path)
        else:
            try:
                text = extract_docx(path)
            except (KeyError, ValueError, zipfile.BadZipFile):
                text = extract_legacy_office(path)
        return text, "indexed" if len(text) >= MIN_TEXT_CHARS else "empty"
    if extension in {".xlsx", ".xlsm"}:
        text = extract_xls(path) if header.startswith(bytes.fromhex("d0cf11e0")) else extract_xlsx(path)
        return text, "indexed" if len(text) >= MIN_TEXT_CHARS else "empty"
    if extension == ".xls":
        text = extract_xlsx(path) if header.startswith(b"PK") else extract_xls(path)
        return text, "indexed" if len(text) >= MIN_TEXT_CHARS else "empty"
    if extension == ".pptx":
        text = extract_pptx(path)
        return text, "indexed" if len(text) >= MIN_TEXT_CHARS else "empty"
    if extension in {".doc", ".ppt", ".wps"}:
        text = extract_legacy_office(path)
        return text, "indexed" if len(text) >= MIN_TEXT_CHARS else "empty_non_content"
    if extension in {
        ".txt",
        ".md",
        ".csv",
        ".json",
        ".jsonl",
        ".html",
        ".xml",
        ".yaml",
        ".yml",
    }:
        text = extract_plain(path)
        return text, "indexed" if len(text) >= MIN_TEXT_CHARS else "empty"
    if header and not header.startswith((b"PK", b"%PDF", bytes.fromhex("d0cf11e0"))):
        raw = path.read_bytes()
        for encoding in ("utf-8", "gb18030"):
            try:
                text = clean_text(raw.decode(encoding))
                return text, "indexed" if len(text) >= MIN_TEXT_CHARS else "empty_non_content"
            except UnicodeDecodeError:
                continue
    return "", "not_text"


def iter_chunks(text: str, size: int = 1_200, overlap: int = 180):
    if not text:
        return
    start = 0
    chunk_number = 1
    while start < len(text):
        end = min(start + size, len(text))
        if end < len(text):
            boundary = max(text.rfind("\n", start, end), text.rfind("。", start, end))
            if boundary > start + size // 2:
                end = boundary + 1
        yield chunk_number, text[start:end]
        if end >= len(text):
            break
        start = max(end - overlap, start + 1)
        chunk_number += 1


def enterprise_mentions(text: str) -> list[tuple[str, str, str]]:
    mentions: list[tuple[str, str, str]] = []
    previous_sequence = ""
    pending_sequence_lines = 0
    seen: set[tuple[str, str]] = set()

    for row_html in HTML_ROW_PATTERN.findall(text):
        cells = [
            html.unescape(HTML_TAG_PATTERN.sub("", cell)).strip()
            for cell in HTML_CELL_PATTERN.findall(row_html)
        ]
        sequence = _normalized_sequence(cells[0]) if cells else ""
        context = " | ".join(cells)[:500]
        for cell in cells:
            current_name_cell = ENTERPRISE_ALIAS_QUALIFIER_PATTERN.sub("", cell)
            for match in ENTERPRISE_PATTERN.finditer(current_name_cell):
                name = match.group(0).strip(" ：:，,、；;。")
                key = (name, sequence)
                if len(name) >= 6 and key not in seen:
                    seen.add(key)
                    mentions.append((name, sequence, context))

    plain_text = HTML_ROW_PATTERN.sub("", text)
    for line in plain_text.splitlines():
        stripped = line.strip()
        standalone_sequence = _normalized_sequence(stripped)
        if standalone_sequence:
            previous_sequence = standalone_sequence
            pending_sequence_lines = 3
            continue
        cells = [cell.strip() for cell in stripped.split("|")]
        sequence = _normalized_sequence(cells[0]) if cells else ""
        sequence = sequence or previous_sequence
        matched_enterprise = False
        current_name_line = ENTERPRISE_ALIAS_QUALIFIER_PATTERN.sub("", stripped)
        for match in ENTERPRISE_PATTERN.finditer(current_name_line):
            name = match.group(0).strip(" ：:，,、；;。")
            name = re.sub(r"^\d+(?:\.0)?[.、\s]*", "", name)
            key = (name, sequence)
            if len(name) < 6 or key in seen:
                continue
            seen.add(key)
            mentions.append((name, sequence, stripped[:500]))
            matched_enterprise = True
        if sequence and not matched_enterprise:
            candidate = re.sub(
                r"^\d+(?:\.0)?[.、\s]*", "", current_name_line
            ).strip(" ：:，,、；;。")
            key = (candidate, sequence)
            if LIST_ENTITY_LINE_PATTERN.fullmatch(candidate) and key not in seen:
                seen.add(key)
                mentions.append((candidate, sequence, stripped[:500]))
                matched_enterprise = True
        if matched_enterprise:
            previous_sequence = ""
            pending_sequence_lines = 0
        elif previous_sequence and stripped:
            pending_sequence_lines -= 1
            if pending_sequence_lines <= 0:
                previous_sequence = ""
    return mentions


def _normalized_sequence(value: str) -> str:
    """Normalize spreadsheet row numbers such as ``368.0`` without treating prose as a sequence."""
    match = re.fullmatch(r"\s*(\d+)(?:\.0+)?\s*", value or "")
    return match.group(1) if match else ""


def structured_small_giant_entities(text: str) -> list[tuple[object, ...]]:
    try:
        payload = json.loads(text)
    except (TypeError, json.JSONDecodeError):
        return []
    if not isinstance(payload, dict) or payload.get("projectId") != 98:
        return []
    if "国家专精特新小巨人" not in str(payload.get("dataset") or ""):
        return []
    records = payload.get("records")
    if not isinstance(records, list):
        return []
    rows: list[tuple[object, ...]] = []
    for sequence, record in enumerate(records, start=1):
        if not isinstance(record, dict):
            continue
        enterprise_name = str(record.get("entName") or "").strip()
        if not enterprise_name:
            continue
        years = sorted({int(year) for year in YEAR_PATTERN.findall(str(record.get("subsidyYear") or ""))})
        earliest_platform_year = years[0] if years else None
        region = "".join(
            str(record.get(field) or "").strip()
            for field in ("province", "city", "county")
        )
        context = json.dumps(
            {
                "平台企业ID": str(record.get("eid") or "").strip(),
                "登记地区": region,
                "行业": str(record.get("industryName") or "").strip(),
                "平台年份": str(record.get("subsidyYear") or "").strip(),
                "最早平台年份": earliest_platform_year,
                "数据状态": "企策顾问动态索引，待逐批官方名单核验",
            },
            ensure_ascii=False,
            separators=(",", ":"),
        )
        rows.append(
            (
                enterprise_name,
                str(sequence),
                "国家专精特新“小巨人”企业",
                None,
                "",
                region,
                "平台历史获批记录待官方名单核验",
                context,
                "medium",
            )
        )
    return rows


def create_database(path: Path, rows: list[dict[str, object]]) -> None:
    catalog = load_project_catalog()
    enriched_rows: list[dict[str, object]] = []
    for row in rows:
        metadata = infer_document_metadata(
            str(row["title"]),
            str(row["source"]),
            str(row["content"]),
            str(row["document_role"]),
            catalog,
        )
        enriched = {**row, **metadata}
        if str(row.get("sensitivity") or "") in {"restricted", "confidential"}:
            enriched["upload_action"] = "restricted_excluded"
        enriched_rows.append(enriched)
    with tempfile.TemporaryDirectory(prefix="jiaotang-kb-content-") as directory:
        temporary_path = Path(directory) / path.name
        connection = sqlite3.connect(temporary_path)
        try:
            connection.executescript(
                """
                CREATE TABLE documents (
                    id INTEGER PRIMARY KEY,
                    source_key TEXT NOT NULL UNIQUE,
                    title TEXT NOT NULL,
                    content TEXT NOT NULL,
                    source TEXT NOT NULL,
                    cloud_path TEXT NOT NULL,
                    document_role TEXT NOT NULL,
                    sensitivity TEXT NOT NULL,
                    sha256 TEXT NOT NULL,
                    updated_at TEXT NOT NULL,
                    canonical_project_name TEXT NOT NULL DEFAULT '',
                    region TEXT NOT NULL DEFAULT '',
                    document_stage TEXT NOT NULL DEFAULT '其他',
                    validity_status TEXT NOT NULL DEFAULT 'active_candidate',
                    policy_year INTEGER,
                    batch TEXT NOT NULL DEFAULT '',
                    replacement_title TEXT NOT NULL DEFAULT '',
                    replacement_basis TEXT NOT NULL DEFAULT '',
                    replacement_url TEXT NOT NULL DEFAULT '',
                    project_id TEXT NOT NULL DEFAULT '',
                    case_pack_id TEXT NOT NULL DEFAULT '',
                    document_type TEXT NOT NULL DEFAULT 'reference_document',
                    evidence_type TEXT NOT NULL DEFAULT '',
                    upload_action TEXT NOT NULL DEFAULT 'review',
                    verification_status TEXT NOT NULL DEFAULT 'unreviewed',
                    parent_document_id INTEGER REFERENCES documents(id),
                    attachment_of INTEGER REFERENCES documents(id),
                    supersedes INTEGER REFERENCES documents(id)
                );
                CREATE VIRTUAL TABLE documents_fts USING fts5(
                    title,
                    content,
                    source,
                    document_role,
                    content='documents',
                    content_rowid='id',
                    tokenize='unicode61'
                );
                CREATE VIRTUAL TABLE documents_fts_trigram USING fts5(
                    title,
                    content,
                    source,
                    document_role,
                    content='documents',
                    content_rowid='id',
                    tokenize='trigram'
                );
                CREATE TABLE document_chunks (
                    id INTEGER PRIMARY KEY,
                    document_id INTEGER NOT NULL REFERENCES documents(id),
                    chunk_number INTEGER NOT NULL,
                    content TEXT NOT NULL,
                    UNIQUE(document_id, chunk_number)
                );
                CREATE VIRTUAL TABLE document_chunks_fts USING fts5(
                    document_id UNINDEXED,
                    chunk_number UNINDEXED,
                    title,
                    content,
                    source,
                    tokenize='trigram'
                );
                CREATE TABLE enterprise_mentions (
                    id INTEGER PRIMARY KEY,
                    document_id INTEGER NOT NULL REFERENCES documents(id),
                    enterprise_name TEXT NOT NULL,
                    sequence_no TEXT NOT NULL,
                    context TEXT NOT NULL,
                    UNIQUE(document_id, enterprise_name, sequence_no)
                );
                CREATE INDEX enterprise_mentions_name_idx
                    ON enterprise_mentions(enterprise_name);
                CREATE TABLE public_list_entities (
                    id INTEGER PRIMARY KEY,
                    document_id INTEGER NOT NULL REFERENCES documents(id),
                    enterprise_name TEXT NOT NULL,
                    sequence_no TEXT NOT NULL,
                    canonical_project_name TEXT NOT NULL,
                    policy_year INTEGER,
                    batch TEXT NOT NULL,
                    region TEXT NOT NULL,
                    list_status TEXT NOT NULL,
                    context TEXT NOT NULL,
                    confidence TEXT NOT NULL,
                    UNIQUE(document_id, enterprise_name, sequence_no)
                );
                CREATE INDEX public_list_entities_name_idx
                    ON public_list_entities(enterprise_name);
                CREATE INDEX public_list_entities_project_idx
                    ON public_list_entities(canonical_project_name, policy_year, region);
                CREATE TABLE public_list_entity_years (
                    id INTEGER PRIMARY KEY,
                    entity_id INTEGER NOT NULL REFERENCES public_list_entities(id) ON DELETE CASCADE,
                    year INTEGER NOT NULL,
                    year_role TEXT NOT NULL,
                    UNIQUE(entity_id, year)
                );
                CREATE INDEX public_list_entity_years_year_idx
                    ON public_list_entity_years(year, entity_id);
                CREATE INDEX documents_policy_metadata_idx
                    ON documents(canonical_project_name, region, document_stage, validity_status);
                CREATE INDEX documents_case_pack_idx
                    ON documents(project_id,case_pack_id,document_type,verification_status);
                CREATE TABLE case_packs (
                    case_pack_id TEXT PRIMARY KEY,
                    project_id TEXT NOT NULL DEFAULT '',
                    project_name TEXT NOT NULL DEFAULT '',
                    title TEXT NOT NULL,
                    enterprise_name TEXT NOT NULL DEFAULT '',
                    year INTEGER,
                    batch TEXT NOT NULL DEFAULT '',
                    industry TEXT NOT NULL DEFAULT '',
                    enterprise_scale TEXT NOT NULL DEFAULT '',
                    sensitivity TEXT NOT NULL DEFAULT 'internal',
                    verification_status TEXT NOT NULL DEFAULT 'auto_grouped',
                    source_root TEXT NOT NULL DEFAULT '',
                    document_count INTEGER NOT NULL DEFAULT 0,
                    created_at TEXT NOT NULL
                );
                CREATE INDEX case_packs_project_idx
                    ON case_packs(project_id,year,verification_status);
                CREATE TABLE case_pack_documents (
                    case_pack_id TEXT NOT NULL REFERENCES case_packs(case_pack_id),
                    document_id INTEGER NOT NULL UNIQUE REFERENCES documents(id),
                    document_type TEXT NOT NULL,
                    evidence_type TEXT NOT NULL DEFAULT '',
                    sequence INTEGER NOT NULL DEFAULT 0,
                    PRIMARY KEY(case_pack_id,document_id)
                );
                CREATE TABLE document_relations (
                    source_document_id INTEGER NOT NULL REFERENCES documents(id),
                    target_document_id INTEGER NOT NULL REFERENCES documents(id),
                    relation_type TEXT NOT NULL,
                    evidence TEXT NOT NULL DEFAULT '',
                    PRIMARY KEY(source_document_id,target_document_id,relation_type)
                );
                CREATE TABLE project_alias_corrections (
                    id INTEGER PRIMARY KEY,
                    raw_project_name TEXT NOT NULL,
                    canonical_project_name TEXT NOT NULL,
                    region TEXT NOT NULL DEFAULT '',
                    start_year INTEGER,
                    end_year INTEGER,
                    status TEXT NOT NULL DEFAULT 'pending',
                    confirmed_by TEXT NOT NULL DEFAULT '',
                    confirmed_at TEXT,
                    note TEXT NOT NULL DEFAULT '',
                    created_at TEXT NOT NULL,
                    updated_at TEXT NOT NULL,
                    UNIQUE(raw_project_name,canonical_project_name,region,start_year,end_year)
                );
                CREATE INDEX project_alias_corrections_lookup_idx
                    ON project_alias_corrections(raw_project_name,status,region);
                CREATE UNIQUE INDEX project_alias_corrections_scope_idx
                    ON project_alias_corrections(
                        raw_project_name,canonical_project_name,region,
                        COALESCE(start_year,0),COALESCE(end_year,9999)
                    );
                CREATE TABLE metadata_match_evidence (
                    id INTEGER PRIMARY KEY,
                    document_id INTEGER NOT NULL REFERENCES documents(id),
                    field_name TEXT NOT NULL,
                    inferred_value TEXT NOT NULL,
                    matched_term TEXT NOT NULL,
                    match_method TEXT NOT NULL,
                    source_scope TEXT NOT NULL,
                    source_excerpt TEXT NOT NULL,
                    rule_version TEXT NOT NULL,
                    confidence TEXT NOT NULL,
                    review_status TEXT NOT NULL DEFAULT 'unreviewed',
                    correction_id INTEGER REFERENCES project_alias_corrections(id),
                    created_at TEXT NOT NULL,
                    UNIQUE(document_id,field_name,rule_version)
                );
                CREATE INDEX metadata_match_evidence_review_idx
                    ON metadata_match_evidence(review_status,confidence,field_name);
                CREATE TABLE policy_verification_queue (
                    id INTEGER PRIMARY KEY,
                    document_id INTEGER NOT NULL REFERENCES documents(id),
                    reason TEXT NOT NULL,
                    priority TEXT NOT NULL,
                    status TEXT NOT NULL DEFAULT 'pending',
                    official_source_url TEXT NOT NULL DEFAULT '',
                    official_document_title TEXT NOT NULL DEFAULT '',
                    official_published_at TEXT,
                    verification_note TEXT NOT NULL DEFAULT '',
                    verified_by TEXT NOT NULL DEFAULT '',
                    verified_at TEXT,
                    created_at TEXT NOT NULL,
                    updated_at TEXT NOT NULL,
                    UNIQUE(document_id,reason)
                );
                CREATE INDEX policy_verification_queue_status_idx
                    ON policy_verification_queue(status,priority,document_id);
                CREATE TABLE policy_document_clusters (
                    id INTEGER PRIMARY KEY,
                    cluster_key TEXT NOT NULL UNIQUE,
                    normalized_title TEXT NOT NULL,
                    document_number TEXT NOT NULL DEFAULT '',
                    canonical_project_name TEXT NOT NULL DEFAULT '',
                    region TEXT NOT NULL DEFAULT '',
                    policy_year INTEGER,
                    representative_document_id INTEGER NOT NULL REFERENCES documents(id),
                    match_method TEXT NOT NULL,
                    confidence TEXT NOT NULL,
                    rule_version TEXT NOT NULL,
                    created_at TEXT NOT NULL,
                    updated_at TEXT NOT NULL
                );
                CREATE INDEX policy_document_clusters_number_idx
                    ON policy_document_clusters(document_number,policy_year,region);
                CREATE TABLE policy_document_cluster_members (
                    id INTEGER PRIMARY KEY,
                    cluster_id INTEGER NOT NULL REFERENCES policy_document_clusters(id),
                    document_id INTEGER NOT NULL UNIQUE REFERENCES documents(id),
                    membership_basis TEXT NOT NULL,
                    confidence TEXT NOT NULL,
                    created_at TEXT NOT NULL
                );
                CREATE INDEX policy_document_cluster_members_cluster_idx
                    ON policy_document_cluster_members(cluster_id,document_id);
                CREATE TABLE policy_verification_propagations (
                    id INTEGER PRIMARY KEY,
                    source_queue_id INTEGER NOT NULL REFERENCES policy_verification_queue(id),
                    cluster_id INTEGER NOT NULL REFERENCES policy_document_clusters(id),
                    source_document_id INTEGER NOT NULL REFERENCES documents(id),
                    target_document_id INTEGER NOT NULL REFERENCES documents(id),
                    field_name TEXT NOT NULL,
                    propagated_value TEXT NOT NULL,
                    official_source_url TEXT NOT NULL DEFAULT '',
                    evidence_excerpt TEXT NOT NULL,
                    rule_version TEXT NOT NULL,
                    propagated_by TEXT NOT NULL,
                    propagated_at TEXT NOT NULL,
                    UNIQUE(source_queue_id,target_document_id,field_name)
                );
                CREATE INDEX policy_verification_propagations_target_idx
                    ON policy_verification_propagations(target_document_id,propagated_at);
                """
            )
            connection.executemany(
                """
                INSERT INTO documents(
                    source_key,title,content,source,cloud_path,document_role,
                    sensitivity,sha256,updated_at,canonical_project_name,region,
                    document_stage,validity_status,policy_year,batch,replacement_title,
                    replacement_basis,replacement_url,project_id,case_pack_id,
                    document_type,evidence_type,upload_action,verification_status
                ) VALUES (?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?)
                """,
                (
                    (
                        row["source_key"],
                        row["title"],
                        row["content"],
                        row["source"],
                        row["cloud_path"],
                        row["document_role"],
                        row["sensitivity"],
                        row["sha256"],
                        row["updated_at"],
                        row["canonical_project_name"],
                        row["region"],
                        row["document_stage"],
                        row["validity_status"],
                        row["policy_year"],
                        row["batch"],
                        row["replacement_title"],
                        row["replacement_basis"],
                        row["replacement_url"],
                        row["project_id"],
                        row["case_pack_id"],
                        row["document_type"],
                        row["evidence_type"],
                        row["upload_action"],
                        row["verification_status"],
                    )
                    for row in enriched_rows
                ),
            )
            connection.execute(
                "INSERT INTO documents_fts(rowid,title,content,source,document_role) "
                "SELECT id,title,content,source,document_role FROM documents"
            )
            connection.execute(
                "INSERT INTO documents_fts_trigram(rowid,title,content,source,document_role) "
                "SELECT id,title,content,source,document_role FROM documents"
            )
            case_rows = [row for row in enriched_rows if row.get("case_pack_id")]
            case_groups: dict[str, list[dict[str, object]]] = {}
            for row in case_rows:
                case_groups.setdefault(str(row["case_pack_id"]), []).append(row)
            for case_pack_id, members in sorted(case_groups.items()):
                first = members[0]
                sensitivity = "restricted" if any(
                    str(item.get("sensitivity")) in {"restricted", "confidential"}
                    for item in members
                ) else "internal"
                connection.execute(
                    """
                    INSERT INTO case_packs(
                        case_pack_id,project_id,project_name,title,year,batch,
                        sensitivity,verification_status,source_root,document_count,created_at
                    ) VALUES (?,?,?,?,?,?,?,?,?,?,?)
                    """,
                    (
                        case_pack_id,
                        str(first.get("project_id") or ""),
                        str(first.get("canonical_project_name") or ""),
                        str(first.get("case_pack_title") or case_pack_id),
                        first.get("policy_year"),
                        str(first.get("batch") or ""),
                        sensitivity,
                        "auto_grouped",
                        str(first.get("case_pack_source_root") or ""),
                        len(members),
                        str(first.get("updated_at") or ""),
                    ),
                )
            connection.execute(
                """
                INSERT INTO case_pack_documents(
                    case_pack_id,document_id,document_type,evidence_type,sequence
                )
                SELECT case_pack_id,id,document_type,evidence_type,
                       ROW_NUMBER() OVER (PARTITION BY case_pack_id ORDER BY id)
                FROM documents WHERE case_pack_id != ''
                """
            )
            connection.execute(
                """
                WITH ranked AS (
                    SELECT case_pack_id,document_id,document_type,
                           FIRST_VALUE(document_id) OVER (
                               PARTITION BY case_pack_id
                               ORDER BY CASE document_type
                                   WHEN 'application' THEN 0
                                   WHEN 'construction_plan' THEN 1
                                   WHEN 'technical_plan' THEN 2
                                   WHEN 'reference_document' THEN 3
                                   ELSE 4 END,
                                   sequence,document_id
                           ) AS root_document_id
                    FROM case_pack_documents
                )
                INSERT INTO document_relations(
                    source_document_id,target_document_id,relation_type,evidence
                )
                SELECT document_id,root_document_id,
                       CASE WHEN document_type='evidence_attachment'
                            THEN 'attachment_of' ELSE 'case_component_of' END,
                       'auto_grouped_case_pack_v1'
                FROM ranked
                WHERE document_id<>root_document_id
                """
            )
            connection.execute(
                """
                UPDATE documents
                SET parent_document_id=(
                    SELECT relations.target_document_id
                    FROM document_relations relations
                    WHERE relations.source_document_id=documents.id
                      AND relations.relation_type IN ('attachment_of','case_component_of')
                    ORDER BY relations.target_document_id LIMIT 1
                ),
                    attachment_of=(
                    SELECT relations.target_document_id
                    FROM document_relations relations
                    WHERE relations.source_document_id=documents.id
                      AND relations.relation_type='attachment_of'
                    ORDER BY relations.target_document_id LIMIT 1
                )
                WHERE EXISTS(
                    SELECT 1 FROM document_relations relations
                    WHERE relations.source_document_id=documents.id
                )
                """
            )
            documents = connection.execute(
                """
                SELECT id,source_key,title,content,source,document_role,canonical_project_name,
                       region,document_stage,policy_year,batch
                FROM documents
                """
            ).fetchall()
            structured_by_source_key = {
                str(row["source_key"]): list(row.get("_structured_entities") or [])
                for row in enriched_rows
                if row.get("_structured_entities")
            }
            chunk_rows: list[tuple[int, int, str, str, str]] = []
            mention_rows: list[tuple[int, str, str, str]] = []
            list_entity_rows: list[tuple[object, ...]] = []
            for (
                document_id,
                source_key,
                title,
                content,
                source,
                document_role,
                canonical_project_name,
                region,
                document_stage,
                policy_year,
                batch,
            ) in documents:
                for chunk_number, chunk in iter_chunks(str(content)):
                    chunk_rows.append(
                        (int(document_id), chunk_number, str(title), chunk, str(source))
                    )
                structured_entities = structured_by_source_key.get(str(source_key), [])
                if not structured_entities:
                    structured_entities = structured_small_giant_entities(str(content))
                if structured_entities:
                    for entity in structured_entities:
                        mention_rows.append(
                            (int(document_id), str(entity[0]), str(entity[1]), str(entity[7]))
                        )
                        list_entity_rows.append((int(document_id), *entity))
                    continue
                for name, sequence, context in enterprise_mentions(str(content)):
                    mention_rows.append((int(document_id), name, sequence, context))
                    if document_role == "50_名单与对标":
                        confidence = "high" if canonical_project_name and document_stage != "其他" else "medium"
                        list_entity_rows.append(
                            (
                                int(document_id),
                                name,
                                sequence,
                                str(canonical_project_name),
                                policy_year,
                                str(batch),
                                str(region),
                                str(document_stage),
                                context,
                                confidence,
                            )
                        )
            connection.executemany(
                "INSERT INTO document_chunks(document_id,chunk_number,content) VALUES (?,?,?)",
                ((row[0], row[1], row[3]) for row in chunk_rows),
            )
            connection.executemany(
                "INSERT INTO document_chunks_fts(document_id,chunk_number,title,content,source) "
                "VALUES (?,?,?,?,?)",
                chunk_rows,
            )
            connection.executemany(
                "INSERT OR IGNORE INTO enterprise_mentions(" 
                "document_id,enterprise_name,sequence_no,context) VALUES (?,?,?,?)",
                mention_rows,
            )
            connection.executemany(
                """
                INSERT OR IGNORE INTO public_list_entities(
                    document_id,enterprise_name,sequence_no,canonical_project_name,
                    policy_year,batch,region,list_status,context,confidence
                ) VALUES (?,?,?,?,?,?,?,?,?,?)
                """,
                list_entity_rows,
            )
            entity_year_rows: list[tuple[int, int, str]] = []
            for entity_id, policy_year, context, confidence in connection.execute(
                "SELECT id,policy_year,context,confidence FROM public_list_entities"
            ):
                years = {int(year) for year in YEAR_PATTERN.findall(str(context or ""))}
                if policy_year:
                    years.add(int(policy_year))
                year_role = "platform_record" if confidence == "medium" else "official_document_year"
                entity_year_rows.extend(
                    (int(entity_id), year, year_role) for year in sorted(years)
                )
            connection.executemany(
                "INSERT OR IGNORE INTO public_list_entity_years(entity_id,year,year_role) VALUES (?,?,?)",
                entity_year_rows,
            )
            enriched_by_source_key = {
                str(row["source_key"]): row for row in enriched_rows
            }
            for document_id, source_key, document_role in connection.execute(
                "SELECT id,source_key,document_role FROM documents"
            ):
                insert_metadata_audit_records(
                    connection,
                    int(document_id),
                    str(document_role),
                    enriched_by_source_key[str(source_key)],
                )
            rebuild_policy_document_clusters(connection)
            try:
                from scripts.build_document_scopes import rebuild_document_scopes
            except ModuleNotFoundError:
                from build_document_scopes import rebuild_document_scopes

            rebuild_document_scopes(connection)
            connection.commit()
        finally:
            connection.close()
        shutil.copy2(temporary_path, path)


def main() -> None:
    args = parse_args()
    manifest_path = args.manifest.expanduser().resolve()
    output = (args.output or manifest_path.parent).expanduser().resolve()
    output.mkdir(parents=True, exist_ok=True)

    cache_path = args.cache.expanduser().resolve()
    cache_path.parent.mkdir(parents=True, exist_ok=True)
    extraction_cache: dict[str, tuple[str, str]] = {}
    if cache_path.exists():
        for line in cache_path.read_text(encoding="utf-8", errors="ignore").splitlines():
            try:
                cached = json.loads(line)
                if cache_status_reusable(str(cached["status"])):
                    extraction_cache[cached["sha256"]] = (cached["text"], cached["status"])
            except (KeyError, TypeError, json.JSONDecodeError):
                continue

    manifest = [json.loads(line) for line in manifest_path.read_text(encoding="utf-8").splitlines()]
    rows: list[dict[str, object]] = []
    report: list[dict[str, object]] = []
    status_counts: Counter[str] = Counter()
    for position, item in enumerate(manifest, start=1):
        if item["upload_action"] not in {"upload", "reference_duplicate"}:
            status = str(item["upload_action"])
            text = ""
        elif item["upload_action"] == "reference_duplicate":
            status = "duplicate_alias"
            text = ""
        elif item["index_mode"] in {"archive_only", "ocr_required"}:
            status = str(item["index_mode"])
            text = ""
        elif item["sha256"] in extraction_cache:
            text, status = extraction_cache[item["sha256"]]
        elif item["index_mode"] == "manual_review":
            try:
                text, status = extract_manual(Path(item["source_path"]))
            except Exception as error:
                text = ""
                status = f"error:{type(error).__name__}"
        else:
            try:
                text, status = extract(Path(item["source_path"]), item["extension"])
            except Exception as error:
                text = ""
                status = f"error:{type(error).__name__}"
            with cache_path.open("a", encoding="utf-8") as cache_target:
                cache_target.write(
                    json.dumps(
                        {"sha256": item["sha256"], "status": status, "text": text},
                        ensure_ascii=False,
                    )
                    + "\n"
                )
        status_counts[status] += 1
        structured_entities: list[tuple[object, ...]] = []
        if status == "indexed" and item.get("extension") == ".json":
            try:
                structured_entities = structured_small_giant_entities(
                    Path(item["source_path"]).read_text(encoding="utf-8", errors="ignore")
                )
            except OSError:
                structured_entities = []
        content_prefix = str(item.get("content_prefix") or "").strip()
        if status == "indexed" and content_prefix:
            text = f"{content_prefix}\n\n{text}".strip()
        report.append(
            {
                "relative_path": item["relative_path"],
                "status": status,
                "text_chars": len(text),
                "sha256": item["sha256"],
            }
        )
        if status == "indexed":
            row = {
                    "source_key": item.get("source_key") or item["sha256"] or item["relative_path"],
                    "title": item.get("title") or item["name"],
                    "content": text,
                    "source": item["relative_path"],
                    "cloud_path": item["cloud_path"],
                    "document_role": item["document_role"],
                    "sensitivity": item["sensitivity"],
                    "sha256": item["sha256"],
                    "updated_at": item["modified_at"],
                }
            if structured_entities:
                row["_structured_entities"] = structured_entities
            rows.append(row)
        if position % 250 == 0:
            print(f"processed={position}/{len(manifest)} indexed={len(rows)}", flush=True)

    with (output / "documents.jsonl").open("w", encoding="utf-8") as target:
        for row in rows:
            target.write(
                json.dumps(
                    {key: value for key, value in row.items() if not key.startswith("_")},
                    ensure_ascii=False,
                )
                + "\n"
            )
    with (output / "extraction_report.csv").open("w", encoding="utf-8-sig", newline="") as target:
        writer = csv.DictWriter(
            target, fieldnames=["relative_path", "status", "text_chars", "sha256"]
        )
        writer.writeheader()
        writer.writerows(report)
    create_database(output / "knowledge_content.sqlite3", rows)
    summary = {
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "manifest_files": len(manifest),
        "indexed_documents": len(rows),
        "status_counts": dict(status_counts),
        "content_characters": sum(len(str(row["content"])) for row in rows),
    }
    (output / "extraction_summary.json").write_text(
        json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    print(json.dumps(summary, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    if len(sys.argv) == 3 and sys.argv[1] == "--extract-pdf-worker":
        sys.stdout.write(extract_pdf(Path(sys.argv[2])))
    else:
        main()
