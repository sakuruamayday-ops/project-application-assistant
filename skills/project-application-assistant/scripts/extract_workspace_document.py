#!/usr/bin/env python3
"""Extract bounded text from a signed workspace document without network access."""

from __future__ import annotations

import argparse
import csv
import io
import json
import posixpath
import re
import sys
import zipfile
from pathlib import Path
from xml.etree import ElementTree

from document_format_detection import DocumentDetection, detect_document


MAX_OUTPUT_CHARACTERS = 500_000
MAX_JSON_OUTPUT_BYTES = 900 * 1024
MAX_ARCHIVE_MEMBER_BYTES = 64 * 1024 * 1024
MAX_ARCHIVE_TOTAL_BYTES = 256 * 1024 * 1024
MAX_PDF_PAGES = 1_000
MAX_PRESENTATION_SLIDES = 1_000
MAX_ARCHIVE_COMPRESSION_RATIO = 2_000
MAX_ODS_COLUMNS = 16_384
MAX_ODS_ROWS = 1_048_576
MAX_ODS_SHEETS = 1_024
MAX_ODS_EXPANDED_CELLS = 1_000_000
MAX_RTF_BINARY_BYTES = 8 * 1024 * 1024


class UnsafeDocumentError(ValueError):
    """The document exceeds a bounded, read-only parsing limit."""


def safe_xml_root(data: bytes) -> ElementTree.Element:
    upper = data.upper()
    if b"<!DOCTYPE" in upper or b"<!ENTITY" in upper:
        raise UnsafeDocumentError("文档 XML 含有不允许的实体或文档类型声明")
    return ElementTree.fromstring(data)


def local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1]


def bounded_join(parts: list[str]) -> tuple[str, bool]:
    text = "".join(parts)
    compact = re.sub(r" {2,}", " ", text)
    compact = re.sub(r"\n{3,}", "\n\n", compact).strip()
    if len(compact) <= MAX_OUTPUT_CHARACTERS:
        return compact, False
    return compact[:MAX_OUTPUT_CHARACTERS].rstrip(), True


def serialize_result(result: dict[str, object]) -> str:
    """Fit stdout under the signed operation byte limit after JSON escaping."""

    serialized = json.dumps(result, ensure_ascii=False, separators=(",", ":"))
    if len(serialized.encode("utf-8")) <= MAX_JSON_OUTPUT_BYTES:
        return serialized
    text = result.get("text")
    if not isinstance(text, str):
        return serialized
    low = 0
    high = len(text)
    fitted = dict(result)
    fitted["truncated"] = True
    while low < high:
        middle = (low + high + 1) // 2
        fitted["text"] = text[:middle].rstrip()
        candidate = json.dumps(fitted, ensure_ascii=False, separators=(",", ":"))
        if len(candidate.encode("utf-8")) <= MAX_JSON_OUTPUT_BYTES:
            low = middle
        else:
            high = middle - 1
    fitted["text"] = text[:low].rstrip()
    return json.dumps(fitted, ensure_ascii=False, separators=(",", ":"))


def safe_archive(path: Path) -> zipfile.ZipFile:
    archive = zipfile.ZipFile(path)
    total = 0
    for info in archive.infolist():
        if info.is_dir():
            continue
        if info.file_size < 0 or info.file_size > MAX_ARCHIVE_MEMBER_BYTES:
            archive.close()
            raise UnsafeDocumentError("文档内部单个文件过大")
        if info.compress_size > 0 and info.file_size / info.compress_size > MAX_ARCHIVE_COMPRESSION_RATIO:
            archive.close()
            raise UnsafeDocumentError("文档压缩比例超出安全上限")
        total += info.file_size
        if total > MAX_ARCHIVE_TOTAL_BYTES:
            archive.close()
            raise UnsafeDocumentError("文档解压后内容过大")
    return archive


def xml_text(data: bytes, paragraph_tags: set[str]) -> str:
    root = safe_xml_root(data)
    parts: list[str] = []
    for element in root.iter():
        tag = local_name(element.tag)
        if tag == "t" and element.text:
            parts.append(element.text)
        elif tag in paragraph_tags:
            parts.append("\n")
        elif tag == "tab":
            parts.append("\t")
        elif tag == "br":
            parts.append("\n")
    return "".join(parts)


def extract_docx(path: Path) -> dict[str, object]:
    with safe_archive(path) as archive:
        names = set(archive.namelist())
        if "word/document.xml" not in names:
            raise ValueError("DOCX 缺少正文结构")
        ordered = ["word/document.xml"]
        ordered.extend(sorted(name for name in names if re.fullmatch(r"word/(?:header|footer)\d+\.xml", name)))
        parts = [xml_text(archive.read(name), {"p", "tr"}) for name in ordered]
    text, truncated = bounded_join([part + "\n" for part in parts])
    return {"kind": "docx", "status": "extracted", "text": text, "truncated": truncated}


def shared_strings(archive: zipfile.ZipFile) -> list[str]:
    try:
        data = archive.read("xl/sharedStrings.xml")
    except KeyError:
        return []
    root = safe_xml_root(data)
    return ["".join(node.text or "" for node in item.iter() if local_name(node.tag) == "t") for item in root]


def workbook_sheets(archive: zipfile.ZipFile) -> list[tuple[str, str]]:
    try:
        workbook = safe_xml_root(archive.read("xl/workbook.xml"))
        relationships = safe_xml_root(archive.read("xl/_rels/workbook.xml.rels"))
    except KeyError:
        return []
    targets = {
        relation.attrib.get("Id", ""): relation.attrib.get("Target", "")
        for relation in relationships
        if local_name(relation.tag) == "Relationship"
    }
    sheets: list[tuple[str, str]] = []
    for index, node in enumerate(item for item in workbook.iter() if local_name(item.tag) == "sheet"):
        relationship_id = next(
            (value for key, value in node.attrib.items() if local_name(key) == "id"),
            "",
        )
        target = targets.get(relationship_id, "")
        if target.startswith("/"):
            worksheet_path = posixpath.normpath(target.lstrip("/"))
        else:
            worksheet_path = posixpath.normpath(posixpath.join("xl", target))
        if not worksheet_path.startswith("xl/worksheets/") or worksheet_path not in archive.namelist():
            raise ValueError("XLSX 工作表关系无效")
        sheets.append((node.attrib.get("name", f"工作表{index + 1}"), worksheet_path))
    return sheets


def workbook_uses_1904_epoch(archive: zipfile.ZipFile) -> bool:
    root = safe_xml_root(archive.read("xl/workbook.xml"))
    properties = next((node for node in root.iter() if local_name(node.tag) == "workbookPr"), None)
    return properties is not None and properties.attrib.get("date1904", "").lower() in {"1", "true"}


def workbook_number_formats(archive: zipfile.ZipFile) -> list[str]:
    try:
        root = safe_xml_root(archive.read("xl/styles.xml"))
    except KeyError:
        return []
    from openpyxl.styles.numbers import BUILTIN_FORMATS

    custom = {
        int(node.attrib["numFmtId"]): node.attrib.get("formatCode", "General")
        for node in root.iter()
        if local_name(node.tag) == "numFmt" and node.attrib.get("numFmtId", "").isdigit()
    }
    cell_formats = next((node for node in root if local_name(node.tag) == "cellXfs"), None)
    if cell_formats is None:
        return []
    return [
        custom.get(int(node.attrib.get("numFmtId", "0")), BUILTIN_FORMATS.get(int(node.attrib.get("numFmtId", "0")), "General"))
        for node in cell_formats
        if local_name(node.tag) == "xf"
    ]


def xlsx_number_text(value_text: str, format_code: str, uses_1904_epoch: bool) -> str:
    value = float(value_text)
    if "%" in format_code:
        match = re.search(r"\.([0#]+)[^%]*%", format_code)
        places = len(match.group(1)) if match else 0
        return f"{value * 100:.{places}f}%"
    from openpyxl.styles.numbers import is_date_format
    from openpyxl.utils.datetime import CALENDAR_MAC_1904, CALENDAR_WINDOWS_1900, from_excel

    if is_date_format(format_code):
        epoch = CALENDAR_MAC_1904 if uses_1904_epoch else CALENDAR_WINDOWS_1900
        converted = from_excel(value, epoch=epoch)
        return converted.isoformat(sep=" ", timespec="seconds")
    return str(int(value)) if value.is_integer() else format(value, ".15g")


def cell_text(
    cell: ElementTree.Element,
    shared: list[str],
    number_formats: list[str],
    uses_1904_epoch: bool,
) -> str:
    kind = cell.attrib.get("t")
    inline = "".join(node.text or "" for node in cell.iter() if local_name(node.tag) == "t")
    if inline:
        return inline
    value = next((node.text or "" for node in cell if local_name(node.tag) == "v"), "")
    if kind == "s" and value.isdigit():
        index = int(value)
        return shared[index] if index < len(shared) else value
    if kind == "b":
        return "TRUE" if value == "1" else "FALSE"
    if kind in {"e", "str"} or value == "":
        return value
    style_index = int(cell.attrib.get("s", "0")) if cell.attrib.get("s", "0").isdigit() else 0
    format_code = number_formats[style_index] if style_index < len(number_formats) else "General"
    try:
        return xlsx_number_text(value, format_code, uses_1904_epoch)
    except (TypeError, ValueError, OverflowError):
        return value


def xlsx_column_index(reference: str, fallback: int) -> int:
    match = re.fullmatch(r"([A-Z]{1,3})\d+", reference.upper())
    if match is None:
        return fallback
    index = 0
    for letter in match.group(1):
        index = index * 26 + ord(letter) - ord("A") + 1
    if index > 16_384:
        raise ValueError("XLSX 单元格列号超出范围")
    return index


def extract_xlsx(path: Path) -> dict[str, object]:
    with safe_archive(path) as archive:
        worksheets = workbook_sheets(archive)
        if not worksheets:
            raise ValueError("XLSX 缺少工作表")
        shared = shared_strings(archive)
        number_formats = workbook_number_formats(archive)
        uses_1904_epoch = workbook_uses_1904_epoch(archive)
        parts: list[str] = []
        for worksheet_name, worksheet_path in worksheets:
            parts.append(f"## {worksheet_name}\n")
            root = safe_xml_root(archive.read(worksheet_path))
            for row in (node for node in root.iter() if local_name(node.tag) == "row"):
                values: list[str] = []
                next_column = 1
                for cell in (item for item in row if local_name(item.tag) == "c"):
                    column = xlsx_column_index(cell.attrib.get("r", ""), next_column)
                    if column < next_column:
                        raise ValueError("XLSX 单元格顺序无效")
                    values.extend([""] * (column - next_column))
                    values.append(cell_text(cell, shared, number_formats, uses_1904_epoch))
                    next_column = column + 1
                parts.append("\t".join(values).rstrip() + "\n")
    text, truncated = bounded_join(parts)
    return {"kind": "xlsx", "status": "extracted", "sheets": len(worksheets), "text": text, "truncated": truncated}


def odf_attribute(element: ElementTree.Element, name: str, default: str = "") -> str:
    return next((value for key, value in element.attrib.items() if local_name(key) == name), default)


def bounded_repeat(element: ElementTree.Element, name: str, maximum: int) -> int:
    raw = odf_attribute(element, name, "1")
    try:
        value = int(raw)
    except ValueError as error:
        raise ValueError(f"ODF {name} 属性无效") from error
    if value < 1 or value > maximum:
        raise UnsafeDocumentError(f"ODF {name} 超出安全上限")
    return value


def odf_cell_text(cell: ElementTree.Element) -> str:
    paragraphs = [
        "".join(paragraph.itertext()).strip()
        for paragraph in cell.iter()
        if local_name(paragraph.tag) in {"p", "h"}
    ]
    visible = "\n".join(value for value in paragraphs if value)
    if visible:
        return visible
    for attribute in ("string-value", "date-value", "time-value", "boolean-value", "value"):
        value = odf_attribute(cell, attribute)
        if value:
            return value
    return ""


def extract_ods(path: Path) -> dict[str, object]:
    with safe_archive(path) as archive:
        try:
            root = safe_xml_root(archive.read("content.xml"))
        except KeyError as error:
            raise ValueError("ODS 缺少 content.xml") from error
        tables = [node for node in root.iter() if local_name(node.tag) == "table"]
        if not tables:
            raise ValueError("ODS 缺少工作表")
        if len(tables) > MAX_ODS_SHEETS:
            raise UnsafeDocumentError("ODS 工作表数量超出安全上限")
        parts: list[str] = []
        output_size = 0
        expanded_cells = 0
        expanded_rows = 0
        truncated = False
        for table_index, table in enumerate(tables):
            name = odf_attribute(table, "name", f"工作表{table_index + 1}")
            heading = f"## {name}\n"
            if output_size + len(heading) > MAX_OUTPUT_CHARACTERS:
                truncated = True
                break
            parts.append(heading)
            output_size += len(heading)
            for row in (node for node in table.iter() if local_name(node.tag) == "table-row"):
                row_repeat = bounded_repeat(row, "number-rows-repeated", MAX_ODS_ROWS)
                values: list[str] = []
                for cell in row:
                    if local_name(cell.tag) not in {"table-cell", "covered-table-cell"}:
                        continue
                    column_repeat = bounded_repeat(cell, "number-columns-repeated", MAX_ODS_COLUMNS)
                    if len(values) + column_repeat > MAX_ODS_COLUMNS:
                        raise UnsafeDocumentError("ODS 单行列数超出安全上限")
                    expanded_cells += column_repeat * row_repeat
                    if expanded_cells > MAX_ODS_EXPANDED_CELLS:
                        raise UnsafeDocumentError("ODS 重复单元格展开超出安全上限")
                    values.extend([odf_cell_text(cell)] * column_repeat)
                expanded_rows += row_repeat
                if expanded_rows > MAX_ODS_ROWS:
                    raise UnsafeDocumentError("ODS 行数超出安全上限")
                line = "\t".join(values).rstrip() + "\n"
                if not line.strip():
                    row_repeat = min(row_repeat, 1)
                for _ in range(row_repeat):
                    if output_size + len(line) > MAX_OUTPUT_CHARACTERS:
                        remaining = MAX_OUTPUT_CHARACTERS - output_size
                        if remaining > 0:
                            parts.append(line[:remaining])
                        truncated = True
                        break
                    parts.append(line)
                    output_size += len(line)
                if truncated:
                    break
            if truncated:
                break
    text, joined_truncated = bounded_join(parts)
    return {
        "kind": "ods",
        "status": "extracted",
        "sheets": len(tables),
        "text": text,
        "truncated": truncated or joined_truncated,
    }


def extract_odt(path: Path) -> dict[str, object]:
    with safe_archive(path) as archive:
        try:
            root = safe_xml_root(archive.read("content.xml"))
        except KeyError as error:
            raise ValueError("ODT 缺少 content.xml") from error
        parts: list[str] = []
        output_size = 0
        truncated = False
        for node in root.iter():
            if local_name(node.tag) not in {"p", "h"}:
                continue
            value = "".join(node.itertext()).strip()
            if value:
                line = value + "\n"
                if output_size + len(line) > MAX_OUTPUT_CHARACTERS:
                    remaining = MAX_OUTPUT_CHARACTERS - output_size
                    if remaining > 0:
                        parts.append(line[:remaining])
                    truncated = True
                    break
                parts.append(line)
                output_size += len(line)
    text, joined_truncated = bounded_join(parts)
    return {"kind": "odt", "status": "extracted", "text": text, "truncated": truncated or joined_truncated}


def xls_number_text(book: object, cell: object) -> str:
    value = float(cell.value)
    format_code = ""
    try:
        xf = book.xf_list[cell.xf_index]
        format_code = book.format_map[xf.format_key].format_str
    except (AttributeError, IndexError, KeyError):
        pass
    if "%" in format_code:
        match = re.search(r"\.([0#]+)[^%]*%", format_code)
        places = len(match.group(1)) if match else 0
        return f"{value * 100:.{places}f}%"
    return str(int(value)) if value.is_integer() else format(value, ".15g")


def xls_cell_text(book: object, cell: object) -> str:
    import xlrd

    if cell.ctype in (xlrd.XL_CELL_EMPTY, xlrd.XL_CELL_BLANK):
        return ""
    if cell.ctype == xlrd.XL_CELL_DATE:
        value = xlrd.xldate_as_datetime(cell.value, book.datemode)
        return value.isoformat(sep=" ", timespec="seconds")
    if cell.ctype == xlrd.XL_CELL_NUMBER:
        return xls_number_text(book, cell)
    if cell.ctype == xlrd.XL_CELL_BOOLEAN:
        return "TRUE" if cell.value else "FALSE"
    if cell.ctype == xlrd.XL_CELL_ERROR:
        return xlrd.error_text_from_code.get(cell.value, "#ERROR")
    return str(cell.value)


def extract_xls(path: Path) -> dict[str, object]:
    try:
        import xlrd
    except ImportError as error:
        raise RuntimeError("内置 XLS 解析组件不可用") from error
    workbook = xlrd.open_workbook(path, on_demand=True, formatting_info=True)
    sheet_count = workbook.nsheets
    parts: list[str] = []
    try:
        for worksheet in workbook.sheets():
            parts.append(f"## {worksheet.name}\n")
            for row_index in range(worksheet.nrows):
                values = [xls_cell_text(workbook, worksheet.cell(row_index, column)) for column in range(worksheet.ncols)]
                parts.append("\t".join(values).rstrip() + "\n")
    finally:
        workbook.release_resources()
    text, truncated = bounded_join(parts)
    return {"kind": "xls", "status": "extracted", "sheets": sheet_count, "text": text, "truncated": truncated}


def extract_pptx(path: Path) -> dict[str, object]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    # 先执行与 Word/Excel 相同的容器上限，避免只按后缀直接解压或解析实体。
    with safe_archive(path) as archive:
        for name in archive.namelist():
            if name.endswith((".xml", ".rels")):
                safe_xml_root(archive.read(name))
    presentation = Presentation(path)
    if len(presentation.slides) > MAX_PRESENTATION_SLIDES:
        raise UnsafeDocumentError(f"演示文稿超过 {MAX_PRESENTATION_SLIDES} 页")
    parts: list[str] = []
    ocr_slides: list[int] = []
    slide_area = presentation.slide_width * presentation.slide_height

    def shape_text(shapes, scale_x: float = 1, scale_y: float = 1) -> tuple[list[str], bool, bool]:
        text: list[str] = []
        images = False
        large_images = False
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                transform = shape._element.xfrm
                child_extent = transform.chExt if transform is not None else None
                group_scale_x = shape.width / child_extent.cx if child_extent is not None and child_extent.cx else 1
                group_scale_y = shape.height / child_extent.cy if child_extent is not None and child_extent.cy else 1
                grouped, group_images, group_large_images = shape_text(
                    shape.shapes, scale_x * group_scale_x, scale_y * group_scale_y,
                )
                text.extend(grouped)
                images = images or group_images
                large_images = large_images or group_large_images
            elif shape.has_text_frame:
                text.append(shape.text_frame.text)
            elif shape.has_table:
                text.extend("\t".join(cell.text for cell in row.cells) for row in shape.table.rows)
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                images = True
                large_images = large_images or shape.width * scale_x * shape.height * scale_y >= slide_area / 2
        return text, images, large_images

    for index, slide in enumerate(presentation.slides, 1):
        text, images, large_images = shape_text(slide.shapes)
        # 标题能读取不代表正文已读完；大幅扫描图仍需识别，小图标不触发整页 OCR。
        if large_images or (images and not any(item.strip() for item in text)):
            ocr_slides.append(index)
        parts.append(f"\n## 第 {index} 页\n" + "\n".join(text) + "\n")
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame
            if notes is not None and notes.text.strip():
                parts.append("\n演讲者备注\n" + notes.text + "\n")
    text, truncated = bounded_join(parts)
    return {
        "kind": "pptx", "slides": len(presentation.slides),
        "status": "needs_ocr" if ocr_slides else "extracted",
        "text": text, "truncated": truncated, "ocr_pages": ocr_slides,
        "action": "ocr_selected_pages" if ocr_slides else "none",
        "message": "已保留可读文字；指定幻灯片仍含需识别的图片，不能把标题当成完整正文。" if ocr_slides else "",
    }


def extract_pdf(path: Path) -> dict[str, object]:
    try:
        import pymupdf
    except ImportError as error:
        raise RuntimeError("内置 PDF 解析组件不可用") from error
    document = pymupdf.open(path)
    if document.needs_pass:
        document.close()
        return {
            "kind": "pdf", "status": "encrypted_document", "text": "", "truncated": False,
            "action": "provide_password_free_copy",
            "message": "PDF 已加密，请提供解除密码后的副本；原文件不会被修改。",
        }
    pages = document.page_count
    if pages > MAX_PDF_PAGES:
        document.close()
        raise UnsafeDocumentError(f"PDF 页数超过 {MAX_PDF_PAGES} 页")
    parts: list[str] = []
    has_text = False
    ocr_pages: list[int] = []
    try:
        for index in range(pages):
            page = document.load_page(index)
            page_text = page.get_text("text", sort=True)
            readable = bool(re.search(r"[A-Za-z0-9\u3400-\u9fff]", page_text))
            if readable:
                has_text = True
            images = page.get_image_info()
            # 扫描正文上方可能有可提取的页眉，不能因为页眉有字就漏掉整页扫描内容。
            large_raster = any((pymupdf.Rect(image["bbox"]) & page.rect).get_area() >= page.rect.get_area() / 2 for image in images)
            # 转曲文字和矢量图没有图片对象，不能在混合 PDF 中被静默跳过。
            unread_visual = not readable and (bool(images) or bool(page.get_drawings()))
            if large_raster or unread_visual:
                ocr_pages.append(index + 1)
            parts.append(f"\n## 第 {index + 1} 页\n{page_text}\n")
    finally:
        document.close()
    text, truncated = bounded_join(parts)
    needs_ocr = bool(ocr_pages) or not has_text
    if not has_text and not ocr_pages:
        ocr_pages = list(range(1, pages + 1))
    return {
        "kind": "pdf",
        "status": "needs_ocr" if needs_ocr else "extracted",
        "pages": pages,
        "text": text if has_text else "",
        "truncated": truncated,
        "ocr_pages": ocr_pages,
        "action": "ocr_selected_pages" if needs_ocr else "none",
        "message": "已保留可提取文字；需要使用已配置的 OCR 补充指定页面，不能把部分正文当成全文。" if needs_ocr else "",
    }


def decode_plain_text(data: bytes) -> str:
    for encoding in ("utf-8-sig", "gb18030"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise ValueError("文本不是 UTF-8 或 GB18030 编码")


def stable_delimiter(text: str, preferred: str | None) -> str | None:
    if preferred is not None:
        return preferred
    lines = [line for line in text[:64_000].splitlines() if line.strip()]
    if len(lines) < 3:
        return None
    for delimiter in ("\t", ",", ";"):
        counts = [line.count(delimiter) for line in lines]
        if counts[0] > 0 and len(set(counts)) == 1:
            return delimiter
    return None


def delimited_text(text: str, preferred: str | None = None) -> tuple[str, str, bool] | None:
    delimiter = stable_delimiter(text, preferred)
    if delimiter is None:
        return None
    output: list[str] = []
    size = 0
    try:
        reader = csv.reader(io.StringIO(text), delimiter=delimiter)
        for row_index, row in enumerate(reader, start=1):
            if row_index > MAX_ODS_ROWS or len(row) > MAX_ODS_COLUMNS:
                raise UnsafeDocumentError("分隔文本行列数超出安全上限")
            line = "\t".join(row).rstrip() + "\n"
            if size + len(line) > MAX_OUTPUT_CHARACTERS:
                remaining = MAX_OUTPUT_CHARACTERS - size
                if remaining > 0:
                    output.append(line[:remaining])
                return "".join(output).rstrip(), "csv" if delimiter != "\t" else "tsv", True
            output.append(line)
            size += len(line)
    except csv.Error:
        return None
    return "".join(output).strip(), "csv" if delimiter != "\t" else "tsv", False


def extract_txt(path: Path) -> dict[str, object]:
    text = decode_plain_text(path.read_bytes())
    suffix = path.suffix.casefold()
    preferred = "\t" if suffix == ".tsv" else ("," if suffix == ".csv" else None)
    table = delimited_text(text, preferred) if suffix in {".csv", ".tsv", ".et"} or stable_delimiter(text, None) else None
    if table is not None:
        table_text, kind, truncated = table
        return {"kind": kind, "status": "extracted", "text": table_text, "truncated": truncated}
    text, truncated = bounded_join([text])
    return {"kind": "txt", "status": "extracted", "text": text, "truncated": truncated}


RTF_DESTINATIONS = {
    "annotation",
    "author",
    "colortbl",
    "comment",
    "datastore",
    "filetbl",
    "fontemb",
    "fontfile",
    "fonttbl",
    "footer",
    "footerf",
    "footerl",
    "footerr",
    "generator",
    "header",
    "headerf",
    "headerl",
    "headerr",
    "info",
    "listoverridetable",
    "listtable",
    "nonshppict",
    "object",
    "pict",
    "revtbl",
    "rsidtbl",
    "shp",
    "shpinst",
    "stylesheet",
    "themedata",
    "xmlnstbl",
}


def rtf_codec(codepage: int) -> str:
    candidate = f"cp{codepage}"
    try:
        b"".decode(candidate)
        return candidate
    except LookupError:
        return "cp1252"


def extract_rtf(path: Path) -> dict[str, object]:
    raw = path.read_bytes()
    source = raw.decode("latin-1")
    if not source.lstrip("\ufeff \t\r\n").lower().startswith("{\\rtf"):
        raise ValueError("RTF 文件头无效")
    # State is copied at every group boundary so hidden destinations never leak.
    state = {"ignorable": False, "ucskip": 1, "codepage": 1252}
    stack: list[dict[str, object]] = []
    output: list[str] = []
    output_size = 0
    fallback_skip = 0
    truncated = False
    index = 0

    def append(value: str) -> None:
        nonlocal output_size, truncated
        if not value or bool(state["ignorable"]) or truncated:
            return
        remaining = MAX_OUTPUT_CHARACTERS - output_size
        if remaining <= 0:
            truncated = True
            return
        if len(value) > remaining:
            output.append(value[:remaining])
            output_size += remaining
            truncated = True
            return
        output.append(value)
        output_size += len(value)

    while index < len(source):
        char = source[index]
        if char == "{":
            stack.append(dict(state))
            index += 1
            continue
        if char == "}":
            if stack:
                state = stack.pop()
            index += 1
            continue
        if char != "\\":
            if char not in "\r\n":
                if fallback_skip:
                    fallback_skip -= 1
                else:
                    append(char)
            index += 1
            continue

        index += 1
        if index >= len(source):
            break
        symbol = source[index]
        if symbol in "{}\\":
            if fallback_skip:
                fallback_skip -= 1
            else:
                append(symbol)
            index += 1
            continue
        if symbol == "'" and index + 2 < len(source):
            token = source[index + 1 : index + 3]
            try:
                decoded = bytes([int(token, 16)]).decode(rtf_codec(int(state["codepage"])), errors="replace")
            except ValueError:
                decoded = ""
            if fallback_skip:
                fallback_skip -= 1
            else:
                append(decoded)
            index += 3
            continue
        if symbol == "*":
            state["ignorable"] = True
            index += 1
            continue
        if not symbol.isalpha():
            if symbol == "~":
                append("\u00a0")
            elif symbol in {"_", "-"}:
                append("-")
            index += 1
            continue

        word_start = index
        while index < len(source) and source[index].isalpha():
            index += 1
        word = source[word_start:index].lower()
        sign = 1
        if index < len(source) and source[index] == "-":
            sign = -1
            index += 1
        number_start = index
        while index < len(source) and source[index].isdigit():
            index += 1
        argument = sign * int(source[number_start:index]) if index > number_start else None
        if index < len(source) and source[index] == " ":
            index += 1

        if word in RTF_DESTINATIONS:
            state["ignorable"] = True
        elif word == "ansicpg" and argument is not None:
            state["codepage"] = argument
        elif word == "uc" and argument is not None:
            state["ucskip"] = max(0, min(argument, 16))
        elif word == "u" and argument is not None:
            codepoint = argument if argument >= 0 else argument + 65_536
            append(chr(codepoint))
            fallback_skip = int(state["ucskip"])
        elif word in {"par", "line", "row"}:
            append("\n")
        elif word in {"tab", "cell"}:
            append("\t")
        elif word == "bin" and argument is not None:
            if argument < 0 or argument > MAX_RTF_BINARY_BYTES or index + argument > len(source):
                raise UnsafeDocumentError("RTF 二进制块超出安全上限或已损坏")
            index += argument

    text, joined_truncated = bounded_join(output)
    return {"kind": "rtf", "status": "extracted", "text": text, "truncated": truncated or joined_truncated}


def outcome(
    path: Path,
    detection: DocumentDetection,
    status: str,
    message: str,
    action: str,
    *,
    retryable: bool = False,
) -> dict[str, object]:
    return {
        "schema_version": "gongchuang-document-extraction/v1",
        "name": path.name,
        **detection.to_dict(),
        "status": status,
        "message": message,
        "action": action,
        "retryable": retryable,
        "text": "",
        "truncated": False,
    }


def non_extractable_outcome(path: Path, detection: DocumentDetection) -> dict[str, object]:
    kind = detection.detected_kind
    if kind == "doc":
        return outcome(
            path,
            detection,
            "conversion_required",
            "检测到旧式二进制 Word 文档。请另存为 DOCX、ODT 或 RTF 后再读取；原文件不会被修改。",
            "convert_to_supported_format",
        )
    if kind == "ppt":
        return outcome(
            path, detection, "conversion_required",
            "检测到旧式二进制演示文稿。请用原应用另存为 PPTX 或 PDF 后继续；原文件不会被修改。",
            "convert_to_supported_format",
        )
    if kind in {"encrypted-office", "encrypted-archive"}:
        return outcome(
            path,
            detection,
            "encrypted_document",
            "文档已加密，当前只读提取器不会请求、保存或尝试破解密码。",
            "provide_password_free_copy",
        )
    if kind.startswith("damaged-"):
        return outcome(
            path,
            detection,
            "damaged_document",
            "文档容器不完整或已损坏，请从原应用重新另存一份有效副本。",
            "replace_with_valid_copy",
        )
    if kind == "unsafe-archive":
        return outcome(
            path,
            detection,
            "unsafe_document",
            detection.detail or "文档容器超出只读嗅探安全上限。",
            "split_or_resave_document",
        )
    if kind == "proprietary-office":
        return outcome(
            path,
            detection,
            "conversion_required",
            "检测到当前内置解析器无法确认的 WPS／ET／DPS 专有文档，请另存为 DOCX、XLSX、PPTX、PDF、ODT、ODS、RTF 或文本后再读取。",
            "convert_to_supported_format",
        )
    if kind == "unreadable":
        return outcome(path, detection, "unavailable", "文档当前不可读取，请检查权限。", "check_permissions")
    return outcome(
        path,
        detection,
        "unsupported_format",
        "未识别到可安全读取的 Office、ODF、PDF、RTF 或文本结构。",
        "provide_supported_format",
    )


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("document", type=Path)
    args = parser.parse_args()
    path = args.document
    try:
        if not path.is_file() or path.is_symlink():
            raise ValueError("输入必须是普通文件")
        detection = detect_document(path)
        handlers = {
            "docx": extract_docx,
            "xlsx": extract_xlsx,
            "pptx": extract_pptx,
            "ods": extract_ods,
            "odt": extract_odt,
            "xls": extract_xls,
            "pdf": extract_pdf,
            "rtf": extract_rtf,
            "text": extract_txt,
        }
        handler = handlers.get(detection.detected_kind)
        if handler is None:
            print(serialize_result(non_extractable_outcome(path, detection)))
            return 0
        result = handler(path)
        result.update({
            "schema_version": "gongchuang-document-extraction/v1",
            "name": path.name,
            **detection.to_dict(),
            "retryable": False,
            "macro_policy": "ignored-read-only" if detection.contains_macros else "not-present",
            "formula_policy": "cached-values-only" if detection.detected_kind in {"xlsx", "xls", "ods"} else "not-applicable",
            "external_links": "not-followed",
        })
        print(serialize_result(result))
        return 0
    except UnsafeDocumentError as error:
        detection = locals().get("detection", DocumentDetection(path.suffix.casefold(), "unknown"))
        result = outcome(path, detection, "unsafe_document", str(error), "split_or_resave_document")
        print(serialize_result(result))
        return 0
    except (zipfile.BadZipFile, ElementTree.ParseError) as error:
        detection = locals().get("detection", DocumentDetection(path.suffix.casefold(), "unknown"))
        result = outcome(path, detection, "damaged_document", str(error), "replace_with_valid_copy")
        print(serialize_result(result))
        return 0
    except ValueError as error:
        if "detection" in locals():
            result = outcome(path, detection, "damaged_document", str(error), "replace_with_valid_copy")
            print(serialize_result(result))
            return 0
        print(serialize_result({
            "schema_version": "gongchuang-document-extraction/v1",
            "name": path.name,
            "status": "rejected",
            "error": str(error),
        }))
        return 2
    except OSError as error:
        detection = locals().get("detection", DocumentDetection(path.suffix.casefold(), "unreadable", detail=str(error)))
        print(serialize_result(outcome(path, detection, "unavailable", str(error), "check_permissions")))
        return 0
    except RuntimeError as error:
        detection = locals().get("detection", DocumentDetection(path.suffix.casefold(), "component-unavailable"))
        print(serialize_result(outcome(path, detection, "component_unavailable", str(error), "update_client")))
        return 0


if __name__ == "__main__":
    sys.exit(main())
