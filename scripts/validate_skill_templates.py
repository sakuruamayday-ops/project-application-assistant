#!/usr/bin/env python3
"""Validate every template carried by the signed skill suite."""

from __future__ import annotations

import argparse
import json
import shutil
import struct
import subprocess
import tempfile
import zipfile
from html.parser import HTMLParser
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


OFFICE_SUFFIXES = {".docx", ".xlsx", ".pptx"}
LEGACY_OFFICE_SUFFIXES = {".doc", ".xls", ".ppt", ".wps", ".et", ".dps"}
SOURCE_TEMPLATE_SUFFIXES = {".html", ".json", ".md"}


class TemplateValidationError(RuntimeError):
    """Raised when a packaged template cannot complete its supported lifecycle."""


class _HTMLStructureParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.start_tags: list[str] = []

    def handle_starttag(self, tag: str, _attrs: list[tuple[str, str | None]]) -> None:
        self.start_tags.append(tag)


def _relative(path: Path, root: Path) -> str:
    return path.relative_to(root).as_posix()


def discover_templates(root: Path) -> tuple[list[Path], list[Path]]:
    """Return Office templates and generator-source templates under ``skills``."""

    skills = root / "skills"
    office: list[Path] = []
    sources: list[Path] = []
    legacy: list[str] = []
    for path in sorted(skills.glob("**/assets/**/*")):
        if not path.is_file():
            continue
        suffix = path.suffix.lower()
        if suffix in LEGACY_OFFICE_SUFFIXES:
            legacy.append(_relative(path, root))
            continue
        if suffix in OFFICE_SUFFIXES:
            office.append(path)
            continue
        template_named = "template" in path.name.lower() or "模板" in path.name
        if template_named and suffix in SOURCE_TEMPLATE_SUFFIXES:
            sources.append(path)
    if legacy:
        raise TemplateValidationError(
            "正式技能模板包含旧式 Office/WPS 格式：" + "、".join(legacy)
        )
    return office, sources


def _validate_docx(path: Path, roundtrip: Path) -> dict[str, object]:
    document = Document(path)
    before = (len(document.paragraphs), len(document.tables), len(document.sections))
    document.save(roundtrip)
    reopened = Document(roundtrip)
    after = (len(reopened.paragraphs), len(reopened.tables), len(reopened.sections))
    if before != after:
        raise TemplateValidationError(f"DOCX 保存后结构数量漂移：{path}")
    return {"paragraphs": before[0], "tables": before[1], "sections": before[2]}


def _validate_xlsx(path: Path, roundtrip: Path) -> dict[str, object]:
    workbook = load_workbook(path, read_only=False, data_only=False, keep_links=False)
    before = list(workbook.sheetnames)
    workbook.save(roundtrip)
    workbook.close()
    reopened = load_workbook(roundtrip, read_only=True, data_only=False, keep_links=False)
    after = list(reopened.sheetnames)
    reopened.close()
    if not before or before != after:
        raise TemplateValidationError(f"XLSX 保存后工作表漂移：{path}")
    return {"worksheets": before}


def _validate_pptx(path: Path, roundtrip: Path) -> dict[str, object]:
    presentation = Presentation(path)
    before = len(presentation.slides)
    presentation.save(roundtrip)
    after = len(Presentation(roundtrip).slides)
    if before != after:
        raise TemplateValidationError(f"PPTX 保存后幻灯片数量漂移：{path}")
    return {"slides": before}


def _validate_source_template(path: Path) -> dict[str, object]:
    text = path.read_text(encoding="utf-8")
    if not text.strip():
        raise TemplateValidationError(f"模板为空：{path}")
    suffix = path.suffix.lower()
    if suffix == ".json":
        parsed = json.loads(text)
        if not isinstance(parsed, (dict, list)):
            raise TemplateValidationError(f"JSON 模板顶层必须是对象或数组：{path}")
        return {"kind": type(parsed).__name__}
    if suffix == ".html":
        parser = _HTMLStructureParser()
        parser.feed(text)
        if "html" not in parser.start_tags or "body" not in parser.start_tags:
            raise TemplateValidationError(f"HTML 模板缺少 html/body：{path}")
        return {"start_tags": len(parser.start_tags)}
    return {"lines": len(text.splitlines())}


def _resolve_soffice(explicit: str | None) -> str:
    if explicit:
        executable = Path(explicit).expanduser().resolve()
        if executable.is_file():
            return str(executable)
        raise TemplateValidationError(f"指定的 LibreOffice 不存在：{executable}")
    detected = shutil.which("soffice") or shutil.which("libreoffice")
    if not detected:
        raise TemplateValidationError("真实渲染要求 LibreOffice，但当前环境未找到 soffice")
    return detected


def _render_office_template(path: Path, soffice: str, work: Path) -> dict[str, object]:
    output = work / "render"
    profile = work / "lo-profile"
    output.mkdir(parents=True)
    profile.mkdir(parents=True)
    command = [
        soffice,
        "--headless",
        f"-env:UserInstallation={profile.as_uri()}",
        "--convert-to",
        "pdf",
        "--outdir",
        str(output),
        str(path),
    ]
    completed = subprocess.run(command, text=True, capture_output=True, timeout=90)
    pdf = output / f"{path.stem}.pdf"
    if completed.returncode or not pdf.is_file() or pdf.stat().st_size == 0:
        detail = completed.stderr.strip() or completed.stdout.strip() or "未生成 PDF"
        raise TemplateValidationError(f"LibreOffice 渲染失败：{path}：{detail}")
    import fitz

    with fitz.open(pdf) as document:
        pages = document.page_count
    if pages < 1:
        raise TemplateValidationError(f"渲染 PDF 没有页面：{path}")
    return {"pdf_bytes": pdf.stat().st_size, "pdf_pages": pages}


def validate_templates(
    root: Path,
    *,
    render: bool = False,
    soffice: str | None = None,
    expected_office_count: int | None = None,
    expected_source_count: int | None = None,
) -> dict[str, object]:
    """Validate discovery, parsing, round-trip writes, and optional real rendering."""

    root = root.resolve()
    office, sources = discover_templates(root)
    if expected_office_count is not None and len(office) != expected_office_count:
        raise TemplateValidationError(
            f"Office 模板数量漂移：{len(office)} != {expected_office_count}"
        )
    if expected_source_count is not None and len(sources) != expected_source_count:
        raise TemplateValidationError(
            f"生成器模板数量漂移：{len(sources)} != {expected_source_count}"
        )
    if not office:
        raise TemplateValidationError("技能套件未发现任何 Office 模板")

    renderer = _resolve_soffice(soffice) if render else None
    results: list[dict[str, object]] = []
    with tempfile.TemporaryDirectory(prefix="jiaotang-template-validation-") as temp:
        temp_root = Path(temp)
        for index, path in enumerate(office):
            suffix = path.suffix.lower()
            item_root = temp_root / str(index)
            item_root.mkdir()
            roundtrip = item_root / f"roundtrip{suffix}"
            if suffix == ".docx":
                detail = _validate_docx(path, roundtrip)
            elif suffix == ".xlsx":
                detail = _validate_xlsx(path, roundtrip)
            else:
                detail = _validate_pptx(path, roundtrip)
            rendered = _render_office_template(path, renderer, item_root) if renderer else None
            results.append(
                {
                    "path": _relative(path, root),
                    "format": suffix.removeprefix("."),
                    "roundtrip": detail,
                    "render": rendered,
                }
            )
        for path in sources:
            results.append(
                {
                    "path": _relative(path, root),
                    "format": path.suffix.lower().removeprefix("."),
                    "source": _validate_source_template(path),
                }
            )

    return {
        "status": "pass",
        "office_template_count": len(office),
        "source_template_count": len(sources),
        "rendered": render,
        "templates": results,
    }


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="检查正式技能套件携带的全部模板")
    parser.add_argument("--root", type=Path, default=Path(__file__).resolve().parents[1])
    parser.add_argument("--render", action="store_true")
    parser.add_argument("--soffice")
    parser.add_argument("--expected-office-count", type=int)
    parser.add_argument("--expected-source-count", type=int)
    parser.add_argument("--output", type=Path)
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    try:
        result = validate_templates(
            args.root,
            render=args.render,
            soffice=args.soffice,
            expected_office_count=args.expected_office_count,
            expected_source_count=args.expected_source_count,
        )
    except Exception as exc:
        result = {"status": "fail", "error": str(exc)}
    payload = json.dumps(result, ensure_ascii=False, indent=2) + "\n"
    if args.output:
        args.output.parent.mkdir(parents=True, exist_ok=True)
        args.output.write_text(payload, encoding="utf-8")
    print(payload, end="")
    return 0 if result["status"] == "pass" else 2


if __name__ == "__main__":
    raise SystemExit(main())
